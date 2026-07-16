#!/usr/bin/env python3
"""湖南省网络安全产业全景图 PPT 生成脚本 — 蓝青色系"""

from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

from ppt_engine import (
    THEME_CYBER, SLIDE_H, WHITE, TEXT_DARK, TEXT_BODY, TEXT_SECONDARY, TEXT_MUTED,
    CONTENT_X, LINE_THIN, LINE_MED, LINE_THICK,
    create_presentation, set_slide_bg, add_rect, add_rounded_rect, add_text,
    add_h_line, add_v_line, page_header,
)

# === 蓝青色系别名 ===
CYBER_DEEP = THEME_CYBER.deep
CYBER_PRIMARY = THEME_CYBER.primary
CYBER_MIST = THEME_CYBER.mist
ACCENT_TEAL = THEME_CYBER.accent_teal
ACCENT_GREEN = THEME_CYBER.accent_green
ACCENT_ORANGE = THEME_CYBER.accent_orange
ACCENT_RED = THEME_CYBER.accent_red
ACCENT_PURPLE = THEME_CYBER.accent_purple
ACCENT_GOLD = THEME_CYBER.accent_gold
BG_WHITE = WHITE
BG_LIGHT = THEME_CYBER.bg_light
TEXT_WHITE = WHITE
LINE_LIGHT = THEME_CYBER.line_light
LINE_BLUE = THEME_CYBER.line_primary

prs, W, H = create_presentation()


TOTAL_PAGES = 9


# ========== 幻灯片 1: 封面 ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

# 左侧蓝色装饰区
add_rect(slide, Inches(0), Inches(0), Inches(4.5), H, fill_color=CYBER_DEEP)

# 蓝色区内的装饰线
add_rect(slide, Inches(0.6), Inches(2.8), Inches(1.2), Inches(0.03), fill_color=BG_WHITE)
add_text(slide, Inches(0.6), Inches(1.5), Inches(3.5), Inches(0.4),
         "HUNAN CYBERSECURITY", font_size=12, color=RGBColor(0xB0, 0xCF, 0xE8),
         font_name='Arial')
add_text(slide, Inches(0.6), Inches(1.85), Inches(3.5), Inches(0.4),
         "INDUSTRY PANORAMA", font_size=12, color=RGBColor(0xB0, 0xCF, 0xE8),
         font_name='Arial')

add_text(slide, Inches(0.6), Inches(5.8), Inches(3.5), Inches(0.3),
         "2026 INDUSTRY REPORT", font_size=10, color=RGBColor(0x90, 0xB8, 0xD8),
         font_name='Arial')

# 右侧白色区 — 主标题
add_text(slide, Inches(5.2), Inches(2.2), Inches(7.5), Inches(0.5),
         "产业研究报告", font_size=14, color=ACCENT_ORANGE, font_name='微软雅黑')

add_text(slide, Inches(5.2), Inches(2.7), Inches(7.5), Inches(1.1),
         "湖南省网络安全\n产业全景图",
         font_size=40, color=CYBER_DEEP, bold=True, font_name='微软雅黑',
         line_spacing=1.3)

# 蓝色短线装饰
add_rect(slide, Inches(5.2), Inches(4.15), Inches(0.8), Inches(0.04), fill_color=CYBER_PRIMARY)

add_text(slide, Inches(5.2), Inches(4.3), Inches(7.5), Inches(0.4),
         "\"两芯一生态\"技术体系  ·  十五五规划关键之年  ·  2026年度产业图谱",
         font_size=12, color=TEXT_SECONDARY, font_name='微软雅黑')

# KPI 数据
kpis = [
    ("8,520", "亿元", "绿色智能计算产业规模（2024）"),
    ("1,400+", "家", "集群骨干企业"),
    ("2,300+", "亿元", "集群2025年产值"),
    ("17", "家", "上市企业"),
]
for i, (val, unit, label) in enumerate(kpis):
    x = Inches(5.2 + i * 1.85)
    add_text(slide, x, Inches(5.2), Inches(1.7), Inches(0.5), val,
             font_size=24, color=CYBER_DEEP, bold=True, font_name='Arial')
    add_text(slide, x + Inches(1.2), Inches(5.35), Inches(0.5), Inches(0.3), unit,
             font_size=11, color=TEXT_MUTED)
    add_text(slide, x, Inches(5.65), Inches(1.7), Inches(0.3), label,
             font_size=9, color=TEXT_SECONDARY)

# 底部信息
add_h_line(slide, Inches(5.2), Inches(6.4), Inches(7), color=LINE_LIGHT, weight=Pt(0.5))
add_text(slide, Inches(5.2), Inches(6.5), Inches(7), Inches(0.3),
         "数据来源：湖南省工信厅 · 长沙市政府 · 湖南日报 · 中国新闻网 · 全国工商联",
         font_size=8, color=TEXT_MUTED)


# ========== 幻灯片 2: 目录 ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
page_header(slide, 2, TOTAL_PAGES, "目  录", THEME_CYBER)

contents = [
    ("01", "产业链全景", "网络安全上中下游全链路解析", CYBER_DEEP),
    ("02", "区域布局", "长株潭核心集聚区+衡阳郴州协同", ACCENT_TEAL),
    ("03", "重点企业", "16家龙头与本土骨干企业", ACCENT_ORANGE),
    ("04", "创新平台", "科研平台与高校力量", ACCENT_GREEN),
    ("05", "政策支撑", "8项核心政策体系", ACCENT_RED),
    ("06", "发展趋势", "四大未来方向", ACCENT_PURPLE),
]

for i, (num, title, desc, color) in enumerate(contents):
    col = i % 2
    row = i // 2
    x = Inches(0.9 + col * 6.0)
    y = Inches(1.4 + row * 1.7)

    add_text(slide, x, y, Inches(1.2), Inches(0.8), num,
             font_size=48, color=color, bold=True, font_name='Arial',
             anchor=MSO_ANCHOR.MIDDLE)
    add_v_line(slide, x + Inches(1.3), y + Inches(0.15), Inches(0.7), color=LINE_LIGHT, weight=Pt(1))
    add_text(slide, x + Inches(1.5), y + Inches(0.1), Inches(4), Inches(0.4), title,
             font_size=18, color=TEXT_DARK, bold=True, font_name='微软雅黑')
    add_text(slide, x + Inches(1.5), y + Inches(0.55), Inches(4), Inches(0.3), desc,
             font_size=11, color=TEXT_SECONDARY, font_name='微软雅黑')


# ========== 幻灯片 3: 产业链全景 ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
page_header(slide, 3, TOTAL_PAGES, "网络安全产业链全景", THEME_CYBER)

chain_data = [
    ("上游", "基础软硬件", CYBER_DEEP, [
        ("CPU/芯片设计", "飞腾CPU、鲲鹏CPU、景嘉微GPU\n国科微SSD控制器、毂梁微DSP", "飞腾 · 鲲鹏 · 景嘉微 · 国科微"),
        ("操作系统", "麒麟OS、麒麟信安OS\nopenEuler信创生态", "麒麟软件 · 麒麟信安"),
        ("密码技术", "国产密码算法、安全芯片\n身份认证基础设施", "商用密码产业基地"),
        ("基础软件", "数据库、中间件、工具软件\n2026年南大通用填补DB空白", "南大通用 · 亚信安全"),
        ("网络设备", "交换机、路由器、光纤通信\n5G+安全网络基础设施", "信维电科 · 中国长城"),
    ]),
    ("中游", "安全产品 · 服务", ACCENT_TEAL, [
        ("网络安全产品", "防火墙、IDS/IPS、态势感知\n终端安全、零信任、云安全", "奇安信 · 深信服 · 360"),
        ("安全服务", "安全运营、渗透测试\n城市安全运营\"长沙模式\"", "奇安信 · 安恒信息"),
        ("信创整机", "国产服务器、PC整机\n信创整机占国产市场65%", "中国长城 · 湘江鲲鹏"),
        ("数据安全", "数据加密、隐私计算\n数据脱敏、数据治理", "科创信息 · 麒麟信安"),
        ("安全集成", "等保测评、系统集成\n行业解决方案", "湘邮科技 · 创智和宇"),
    ]),
    ("下游", "行业应用", ACCENT_ORANGE, [
        ("党政信创", "电子政务安全、国产替代\n党政机关安全办公", "麒麟 · 长城 · 奇安信"),
        ("金融安全", "金融信创、安全交易\n金融行业等保合规", "麒麟信安 · 长沙银行"),
        ("能源电力", "电力监控系统安全\n麒麟信安电力行业深耕", "麒麟信安 · 国网湖南"),
        ("电信与算力", "5G安全、算力网络安全\n长株潭+东江湖算力集群", "中国电信 · 曙光云"),
        ("医疗教育", "医疗信创平台、教育信息化\n移动医疗安全", "中南大学湘雅 · 创智和宇"),
    ]),
]

y_start = Inches(1.15)
row_h = Inches(1.88)
for row_idx, (label, sub, color, items) in enumerate(chain_data):
    y = y_start + row_idx * row_h
    add_rect(slide, Inches(0.7), y, Inches(1.3), Inches(1.65), fill_color=color)
    add_text(slide, Inches(0.7), y + Inches(0.3), Inches(1.3), Inches(0.5), label,
             font_size=20, color=TEXT_WHITE, bold=True, alignment=PP_ALIGN.CENTER,
             font_name='微软雅黑')
    add_text(slide, Inches(0.7), y + Inches(0.85), Inches(1.3), Inches(0.3), sub,
             font_size=9, color=RGBColor(0xD0, 0xE0, 0xF0), alignment=PP_ALIGN.CENTER,
             font_name='微软雅黑')

    card_w = Inches(2.22)
    card_gap = Inches(0.08)
    for i, (title, desc, ent) in enumerate(items):
        cx = Inches(2.15) + i * (card_w + card_gap)
        add_rect(slide, cx, y, card_w, Inches(1.65), fill_color=BG_LIGHT,
                 line_color=LINE_LIGHT, line_width=Pt(0.5))
        add_rect(slide, cx, y, card_w, Inches(0.03), fill_color=color)
        add_text(slide, cx + Inches(0.1), y + Inches(0.15), card_w - Inches(0.2), Inches(0.3),
                 title, font_size=12, color=color, bold=True, font_name='微软雅黑')
        add_text(slide, cx + Inches(0.1), y + Inches(0.5), card_w - Inches(0.2), Inches(0.6),
                 desc, font_size=9, color=TEXT_BODY, font_name='微软雅黑', line_spacing=1.3)
        add_h_line(slide, cx + Inches(0.1), y + Inches(1.15), card_w - Inches(0.2),
                   color=LINE_LIGHT, weight=Pt(0.5))
        add_text(slide, cx + Inches(0.1), y + Inches(1.22), card_w - Inches(0.2), Inches(0.35),
                 ent, font_size=8, color=color, font_name='微软雅黑')


# ========== 幻灯片 4: 区域布局 ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
page_header(slide, 4, TOTAL_PAGES, "区域布局与产业集聚区", THEME_CYBER)

regions = [
    ("长沙", "国家网安产业园（全国第二）", "信创整机65%国产市场 · 自主安全计算集群", "产值2300亿+", CYBER_DEEP,
     "奇安信 · 深信服 · 中国长城 · 湘江鲲鹏 · 麒麟软件 · 麒麟信安 · 国科微 · 景嘉微 · 拓维信息 · 科创信息"),
    ("株洲", "功率半导体+智能终端", "中低压功率器件 · 智能终端制造基地", "220+企业", ACCENT_TEAL,
     "中车时代半导体 · 三一硅能 · 赛德雷特"),
    ("湘潭", "产业转化基地", "半导体制造 · 智能装备 · 产学研转化", "23.8亿投资", ACCENT_ORANGE,
     "湘潭智造基地 · 锦智光电 · 蓝思智能终端"),
    ("衡阳", "存储与智能终端", "白沙洲工业园 · 存储产业差异化定位", "衡州大道数字经济带", ACCENT_GREEN,
     "衡阳白沙洲 · 智能终端制造 · 数字经济走廊"),
    ("郴州", "东江湖算力集群", "郴州东江湖大数据中心 · 算力节点", "省级算力集群", ACCENT_RED,
     "东江湖数据中心 · 湾区产业承接"),
    ("益阳", "电子元器件配套", "高端电子元器件 · 光通信器件 · 线路板", "86项专利", ACCENT_PURPLE,
     "信维电科 · 光智通信 · 江丰电子 · 明正宏"),
]

card_w = Inches(4.0)
card_h = Inches(2.65)
gap_x = Inches(0.13)
gap_y = Inches(0.18)
start_x = Inches(0.7)
start_y = Inches(1.15)

for i, (city, badge, role, data_str, color, enterprises) in enumerate(regions):
    col = i % 3
    row = i // 3
    x = start_x + col * (card_w + gap_x)
    y = start_y + row * (card_h + gap_y)

    add_rect(slide, x, y, card_w, card_h, fill_color=BG_WHITE,
             line_color=LINE_LIGHT, line_width=Pt(0.5))
    add_rect(slide, x, y, Inches(0.06), card_h, fill_color=color)
    add_text(slide, x + Inches(0.2), y + Inches(0.12), Inches(2.5), Inches(0.25), badge,
             font_size=9, color=color, bold=True, font_name='微软雅黑')
    add_text(slide, x + Inches(0.2), y + Inches(0.38), Inches(2), Inches(0.45), city,
             font_size=26, color=TEXT_DARK, bold=True, font_name='微软雅黑')
    add_text(slide, x + Inches(2.5), y + Inches(0.45), card_w - Inches(2.6), Inches(0.35), data_str,
             font_size=13, color=color, bold=True, alignment=PP_ALIGN.RIGHT, font_name='微软雅黑')
    add_text(slide, x + Inches(0.2), y + Inches(0.95), card_w - Inches(0.3), Inches(0.3), role,
             font_size=10, color=TEXT_SECONDARY, font_name='微软雅黑')
    add_h_line(slide, x + Inches(0.2), y + Inches(1.3), card_w - Inches(0.4),
               color=LINE_LIGHT, weight=Pt(0.5))
    add_text(slide, x + Inches(0.2), y + Inches(1.4), card_w - Inches(0.3), Inches(1.1), enterprises,
             font_size=9, color=TEXT_BODY, font_name='微软雅黑', line_spacing=1.5)


# ========== 幻灯片 5: 重点企业 ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
page_header(slide, 5, TOTAL_PAGES, "龙头与重点企业", THEME_CYBER)

enterprises = [
    ("奇安信", "长沙·国家网安产业园", "湖南本地最大网安企业，累计投入超8亿元，5亿元长沙奇安基金，Q-GPT安全机器人国内第一", "湖南最大", ACCENT_RED),
    ("深信服", "长沙·研发中心", "深信服长沙研发中心，网络安全+云计算+AI，深度参与长沙网安产业园建设", "巨头布局", ACCENT_RED),
    ("中国长城", "长沙", "国产整机龙头，PK体系核心厂商，信创整机国产市场市占率领先", "央企龙头", ACCENT_RED),
    ("麒麟软件", "长沙·分公司", "国产操作系统龙头（天津总部），麒麟OS为核心，\"两芯一生态\"体系基石，长沙520人", "\"两芯一生态\"", ACCENT_RED),
    ("湘江鲲鹏", "长沙", "华为鲲鹏生态核心伙伴（拓维信息控股90%），国产服务器/PC整机", "鲲鹏生态", ACCENT_ORANGE),
    ("麒麟信安", "长沙·高新区", "科创板上市(688152)，2025营收3.1亿，操作系统+信息安全+云计算三位一体", "上市企业", ACCENT_GOLD),
    ("国科微", "长沙", "国产SSD控制器芯片龙头，\"七大类芯片\"之一，国家级专精特新小巨人", "芯片龙头", ACCENT_TEAL),
    ("景嘉微", "长沙", "国产GPU龙头，军民融合，\"两芯一生态\"体系GPU芯片唯一提供商", "芯片龙头", ACCENT_TEAL),
    ("拓维信息", "长沙", "华为鲲鹏/昇腾生态核心ISV，智慧教育+智慧交通+信创安全", "上市企业", ACCENT_GREEN),
    ("科创信息", "长沙", "智慧政务+数字政府龙头，国产大数据平台，信创安全解决方案", "本土龙头", ACCENT_GREEN),
    ("湘邮科技", "长沙", "邮政行业信息化龙头，信创安全集成服务，数据安全合规", "本土龙头", ACCENT_GREEN),
    ("飞腾信息", "长沙·研发中心", "国产CPU龙头（天津总部），飞腾CPU为\"两芯一生态\"核心之一", "\"两芯一生态\"", ACCENT_ORANGE),
    ("亚信安全", "长沙", "身份安全+终端安全龙头，湖南信创安全生态核心参与者", "巨头布局", ACCENT_PURPLE),
    ("360安全", "长沙", "数字安全服务商，长沙网安产业园入驻企业，政企安全业务", "巨头布局", ACCENT_PURPLE),
    ("创智和宇", "长沙", "医疗信息化龙头，医疗信创安全解决方案，国产替代先行者", "本土龙头", ACCENT_TEAL),
    ("安恒信息", "长沙·研发中心", "态势感知+数据安全，长沙研发与服务团队，参与网安产业园建设", "巨头布局", ACCENT_PURPLE),
]

card_w = Inches(3.05)
card_h = Inches(1.35)
gap_x = Inches(0.1)
gap_y = Inches(0.1)
start_x = Inches(0.7)
start_y = Inches(1.1)

for i, (name, loc, desc, badge, badge_color) in enumerate(enterprises):
    col = i % 4
    row = i // 4
    x = start_x + col * (card_w + gap_x)
    y = start_y + row * (card_h + gap_y)

    add_rect(slide, x, y, card_w, card_h, fill_color=BG_WHITE,
             line_color=LINE_LIGHT, line_width=Pt(0.5))
    add_rect(slide, x, y, Inches(0.04), card_h, fill_color=badge_color if badge else LINE_LIGHT)
    add_text(slide, x + Inches(0.12), y + Inches(0.06), card_w - Inches(0.2), Inches(0.3),
             name, font_size=13, color=TEXT_DARK, bold=True, font_name='微软雅黑')
    add_text(slide, x + Inches(0.12), y + Inches(0.36), card_w - Inches(0.2), Inches(0.2),
             loc, font_size=8, color=TEXT_MUTED, font_name='微软雅黑')
    add_text(slide, x + Inches(0.12), y + Inches(0.58), card_w - Inches(0.2), Inches(0.55),
             desc, font_size=8, color=TEXT_BODY, font_name='微软雅黑', line_spacing=1.3)
    if badge:
        add_text(slide, x + Inches(0.12), y + Inches(1.1), Inches(1.2), Inches(0.2), badge,
                 font_size=8, color=badge_color, bold=True, font_name='微软雅黑')


# ========== 幻灯片 6: 创新平台 ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
page_header(slide, 6, TOTAL_PAGES, "创新平台与科研支撑", THEME_CYBER)

platforms = [
    ("国家网络安全产业园区（长沙）", "全国第二个获批，\"一院四中心\"公共服务平台\n信创适配中心·安全运营中心·测试认证中心·工业互联网推广中心", CYBER_DEEP),
    ("湘江实验室", "湖南四大实验室之一\n聚焦先进计算+AI+网络安全，产学研协同创新", ACCENT_TEAL),
    ("湖南大学网络空间安全学院", "2024年入选国家一流网络安全学院（全国16所）\n依托国家超算长沙中心，大数据攻防试验场", ACCENT_ORANGE),
    ("高校科研力量", "国防科大 · 湖南大学 · 中南大学 · 湖南师大\n湘潭大学 · 湖南科大 · 南华大学（首批一流网安院系6所）", ACCENT_GREEN),
    ("首批省级网信研究基地（7家）", "中南大学·湖南师大·湘潭大学·湘江实验室\n北京大学长沙院·新华社国重/云目未来·方滨兴院士站/文盾信息", ACCENT_RED),
    ("信创人才产教融合基地", "联合国防科大等13所高校+长城等15家企业\n已培养超2万名人才，85%在集群企业就业", ACCENT_PURPLE),
    ("国家网安教育技术产业融合试验区", "2022年湘江新区获批全国首批\n政校企研融合，20+高校全链条人才培育", ACCENT_RED),
    ("湖南省网信研究基地", "2024年1月授牌，7家单位\n32选7严格遴选，网信特色智库", ACCENT_GOLD),
    ("飞地园区 · 科创飞地", "长株潭飞地园区机制\nGDP核算、税收分成、能耗指标突破", CYBER_PRIMARY),
]

card_w = Inches(4.0)
card_h = Inches(1.68)
gap_x = Inches(0.13)
gap_y = Inches(0.15)
start_x = Inches(0.7)
start_y = Inches(1.15)

for i, (title, desc, color) in enumerate(platforms):
    col = i % 3
    row = i // 3
    x = start_x + col * (card_w + gap_x)
    y = start_y + row * (card_h + gap_y)

    add_rect(slide, x, y, card_w, card_h, fill_color=BG_LIGHT,
             line_color=LINE_LIGHT, line_width=Pt(0.5))
    add_rounded_rect(slide, x + Inches(0.15), y + Inches(0.15), Inches(0.4), Inches(0.4),
                    fill_color=color, radius=0.5)
    add_text(slide, x + Inches(0.15), y + Inches(0.15), Inches(0.4), Inches(0.4), str(i+1),
             font_size=14, color=TEXT_WHITE, bold=True, alignment=PP_ALIGN.CENTER,
             font_name='Arial', anchor=MSO_ANCHOR.MIDDLE)
    add_text(slide, x + Inches(0.65), y + Inches(0.15), card_w - Inches(0.8), Inches(0.4), title,
             font_size=12, color=TEXT_DARK, bold=True, font_name='微软雅黑',
             anchor=MSO_ANCHOR.MIDDLE)
    add_text(slide, x + Inches(0.15), y + Inches(0.7), card_w - Inches(0.3), Inches(0.9), desc,
             font_size=9, color=TEXT_BODY, font_name='微软雅黑', line_spacing=1.4)


# ========== 幻灯片 7: 政策支撑 ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
page_header(slide, 7, TOTAL_PAGES, "政策支撑体系", THEME_CYBER)

policies = [
    ("立法保障", "《湖南省网络安全和信息化条例》", "全国率先立法，要求县级以上政府设网信专项资金，引导社会资金投资网信建设", CYBER_DEEP),
    ("产业体系", "\"4×4\"现代化产业体系", "自主安全计算列为主攻方向之首，2025年产业规模突破1万亿，2030年达2万亿", ACCENT_TEAL),
    ("园区政策", "长沙市加快网络安全产业发展政策", "设专项资金+50亿产业基金，龙头企业最高5000万补助，上市奖200万，研发补贴30%", ACCENT_ORANGE),
    ("人才政策", "网安人才实训基地+领军团队", "实训基地30万/年运营补助，高端人才最高200万奖励+全额购房补贴，领军团队最高1亿", ACCENT_GREEN),
    ("区域协同", "长株潭一体化发展行动计划", "\"一核两区三城\"布局，研发在长沙、转化在湘潭、应用在株洲", ACCENT_RED),
    ("十五五规划", "湖南省十五五网信规划", "2025年5月公开征集意见，网络安全保障体系、AI安全风险防范、壮大网安产业", ACCENT_PURPLE),
    ("创新激励", "\"先用后付\"成果转化+适配补助", "537项科技成果推广，信创适配中心每年新增最高200万适配补助", CYBER_PRIMARY),
    ("飞地机制", "飞地园区发展若干措施", "GDP核算、税收分成、能耗指标突破，支持长株潭产业跨区域协同布局", ACCENT_GOLD),
]

card_w = Inches(6.1)
card_h = Inches(1.32)
gap_x = Inches(0.2)
gap_y = Inches(0.1)
start_x = Inches(0.7)
start_y = Inches(1.15)

for i, (tag, title, desc, color) in enumerate(policies):
    col = i % 2
    row = i // 2
    x = start_x + col * (card_w + gap_x)
    y = start_y + row * (card_h + gap_y)

    add_rect(slide, x, y, card_w, card_h, fill_color=BG_WHITE,
             line_color=LINE_LIGHT, line_width=Pt(0.5))
    add_rect(slide, x, y, Inches(0.06), card_h, fill_color=color)
    add_text(slide, x + Inches(0.2), y + Inches(0.1), Inches(0.8), Inches(0.25), tag,
             font_size=9, color=color, bold=True, font_name='微软雅黑')
    add_text(slide, x + Inches(1.1), y + Inches(0.08), card_w - Inches(1.2), Inches(0.35), title,
             font_size=13, color=TEXT_DARK, bold=True, font_name='微软雅黑',
             anchor=MSO_ANCHOR.MIDDLE)
    add_text(slide, x + Inches(0.2), y + Inches(0.55), card_w - Inches(0.3), Inches(0.65), desc,
             font_size=10, color=TEXT_BODY, font_name='微软雅黑', line_spacing=1.3)


# ========== 幻灯片 8: 发展趋势 ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
page_header(slide, 8, TOTAL_PAGES, "发展趋势与未来方向", THEME_CYBER)

trends = [
    ("01", "AI + 安全融合", CYBER_DEEP,
     "奇安信Q-GPT安全机器人国内大模型评测第一，威胁减少80%。湖南率先出台AI安全治理指引，建立AI安全风险监测预警机制。\"十五五\"推动AI+安全升级为\"湖南模式2.0\""),
    ("02", "信创国产替代", ACCENT_TEAL,
     "\"两芯一生态\"体系覆盖4.8万家企业560万款产品，信创整机占国产市场65%。\"六机七芯\"占据市场主导，从党政信创向金融、电力、医疗等8大行业加速渗透"),
    ("03", "数据安全治理", ACCENT_ORANGE,
     "隐私计算+数据脱敏技术突破，曙光云\"立体密算\"理论体系。长沙数据标注基地+湖南大数据交易所，多方协同数据安全治理机制建设"),
    ("04", "长株潭一体化", ACCENT_RED,
     "\"一核两区三城\"空间格局，三市集群梯次培育体系。长株潭+东江湖两大算力集群争取纳入全国一体化算力网络，\"总部在长沙、生产在株潭\"飞地模式深化"),
]

card_w = Inches(2.95)
card_h = Inches(4.8)
gap = Inches(0.15)
start_x = Inches(0.7)
start_y = Inches(1.3)

for i, (num, title, color, text) in enumerate(trends):
    x = start_x + i * (card_w + gap)
    add_rect(slide, x, start_y, card_w, card_h, fill_color=BG_LIGHT,
             line_color=LINE_LIGHT, line_width=Pt(0.5))
    add_rect(slide, x, start_y, card_w, Inches(0.06), fill_color=color)
    add_text(slide, x + Inches(0.2), start_y + Inches(0.25), Inches(2), Inches(1.2), num,
             font_size=56, color=CYBER_MIST, bold=True, font_name='Arial')
    add_rect(slide, x + Inches(0.2), start_y + Inches(1.55), Inches(0.6), Inches(0.03),
             fill_color=color)
    add_text(slide, x + Inches(0.2), start_y + Inches(1.75), card_w - Inches(0.4), Inches(0.5), title,
             font_size=18, color=color, bold=True, font_name='微软雅黑')
    add_text(slide, x + Inches(0.2), start_y + Inches(2.5), card_w - Inches(0.4), Inches(2), text,
             font_size=10, color=TEXT_BODY, font_name='微软雅黑', line_spacing=1.6)


# ========== 幻灯片 9: 总结 ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_rect(slide, Inches(0), Inches(0), Inches(0.35), H, fill_color=CYBER_DEEP)

add_text(slide, Inches(0.8), Inches(0.8), Inches(10), Inches(0.5),
         "总结与展望", font_size=28, color=CYBER_DEEP, bold=True, font_name='微软雅黑')
add_h_line(slide, Inches(0.8), Inches(1.4), Inches(11.8), color=LINE_LIGHT, weight=Pt(1))

add_text(slide, Inches(0.8), Inches(1.7), Inches(11.5), Inches(0.8),
         "湖南网络安全产业 — 迈向自主安全计算新高地",
         font_size=24, color=TEXT_DARK, bold=True, font_name='微软雅黑',
         alignment=PP_ALIGN.CENTER)

summary_items = [
    ("产业规模", "绿色智能计算产业2024年8520亿元（+19.2%），2025年突破1万亿。集群产值超2300亿，全国唯一计算领域国家级集群", CYBER_DEEP),
    ("区域格局", "长沙国家网安产业园（全国第二），长株潭\"一核两区三城\"，衡阳郴州益阳协同，六大集聚区差异化发展", ACCENT_TEAL),
    ("技术体系", "\"两芯一生态\"覆盖4.8万家企业560万款产品，信创整机占国产65%，\"六机七芯\"主导市场", ACCENT_ORANGE),
    ("政策护航", "全国率先网信立法、50亿产业基金、\"十五五\"网信规划、飞地园区机制、\"先用后付\"成果转化", ACCENT_RED),
    ("未来方向", "AI+安全融合、信创国产替代加速、数据安全治理深化、长株潭一体化升级", ACCENT_PURPLE),
]

for i, (title, desc, color) in enumerate(summary_items):
    y = Inches(2.9 + i * 0.72)
    add_rect(slide, Inches(1.0), y, Inches(11), Inches(0.6), fill_color=BG_LIGHT,
             line_color=LINE_LIGHT, line_width=Pt(0.5))
    add_rect(slide, Inches(1.0), y, Inches(0.06), Inches(0.6), fill_color=color)
    add_text(slide, Inches(1.2), y, Inches(2), Inches(0.6), title,
             font_size=14, color=color, bold=True, font_name='微软雅黑',
             anchor=MSO_ANCHOR.MIDDLE)
    add_text(slide, Inches(3.3), y, Inches(8.5), Inches(0.6), desc,
             font_size=11, color=TEXT_BODY, font_name='微软雅黑',
             anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.4)

add_h_line(slide, Inches(0.8), Inches(6.6), Inches(11.5), color=LINE_LIGHT, weight=Pt(0.5))
add_text(slide, Inches(0.8), Inches(6.7), Inches(11.5), Inches(0.3),
         "数据来源：湖南省工业和信息化厅 · 长沙市人民政府 · 湖南日报 · 中国新闻网 · 全国工商联  |  2026年度",
         font_size=9, color=TEXT_MUTED, alignment=PP_ALIGN.CENTER, font_name='微软雅黑')


# === 保存 ===
output_path = "/Users/mrlin/Desktop/qingruan/湖南省网络安全产业全景图.pptx"
prs.save(output_path)
print(f"PPT 已保存至: {output_path}")
print(f"共 {len(prs.slides)} 页幻灯片")
