#!/usr/bin/env python3
"""湖南省半导体行业协会 PPT 生成脚本"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR

# === 红金色系（半导体主题）===
SEMI_DEEP = RGBColor(0x8B, 0x1A, 0x1A)
SEMI_PRIMARY = RGBColor(0xB5, 0x2E, 0x2E)
SEMI_LIGHT = RGBColor(0xD4, 0x4A, 0x3A)
SEMI_PALE = RGBColor(0xFC, 0xF5, 0xF3)
SEMI_MIST = RGBColor(0xF8, 0xEA, 0xE8)

ACCENT_TEAL = RGBColor(0x00, 0x6C, 0x67)
ACCENT_GOLD = RGBColor(0xC4, 0x9A, 0x2A)
ACCENT_ORANGE = RGBColor(0xD4, 0x6A, 0x1A)
ACCENT_NAVY = RGBColor(0x1A, 0x2E, 0x5C)
ACCENT_GREEN = RGBColor(0x2E, 0x86, 0x4B)
ACCENT_RED = RGBColor(0xC0, 0x39, 0x2B)
ACCENT_PURPLE = RGBColor(0x6C, 0x34, 0x80)

BG_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BG_LIGHT = RGBColor(0xFD, 0xF7, 0xF5)
BG_SIDEBAR = SEMI_DEEP

TEXT_DARK = RGBColor(0x2C, 0x2C, 0x2C)
TEXT_BODY = RGBColor(0x4A, 0x4A, 0x4A)
TEXT_SECONDARY = RGBColor(0x7A, 0x7A, 0x8A)
TEXT_MUTED = RGBColor(0xAA, 0xAA, 0xB8)
TEXT_WHITE = RGBColor(0xFF, 0xFF, 0xFF)

LINE_LIGHT = RGBColor(0xE8, 0xDE, 0xDA)
LINE_RED = RGBColor(0xB5, 0x2E, 0x2E)

CAT_COLORS = [SEMI_PRIMARY, ACCENT_ORANGE, ACCENT_GOLD,
              ACCENT_TEAL, ACCENT_NAVY, ACCENT_GREEN, ACCENT_PURPLE]

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
W = prs.slide_width
H = prs.slide_height


def set_slide_bg(slide, color=BG_WHITE):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, fill_color=None, line_color=None, line_width=Pt(0.5)):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.shadow.inherit = False
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape


def add_rounded_rect(slide, left, top, width, height, fill_color=None, line_color=None,
                     line_width=Pt(0.5), radius=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.shadow.inherit = False
    if radius is not None:
        shape.adjustments[0] = radius
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, left, top, width, height, text, font_size=12, color=TEXT_DARK,
             bold=False, alignment=PP_ALIGN.LEFT, font_name='微软雅黑', anchor=MSO_ANCHOR.TOP,
             line_spacing=1.2):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Pt(2)
    tf.margin_right = Pt(2)
    tf.margin_top = Pt(1)
    tf.margin_bottom = Pt(1)
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    p.space_after = Pt(2)
    p.line_spacing = line_spacing
    return txBox


def add_multiline_text(slide, left, top, width, height, lines, font_size=12, color=TEXT_DARK,
                       bold=False, alignment=PP_ALIGN.LEFT, font_name='微软雅黑', line_spacing=1.3):
    """支持多行的文本框"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(2)
    tf.margin_right = Pt(2)
    tf.margin_top = Pt(1)
    tf.margin_bottom = Pt(1)

    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = font_name
        p.alignment = alignment
        p.space_after = Pt(2)
        p.line_spacing = line_spacing
    return txBox


def add_h_line(slide, left, top, width, color=LINE_RED, weight=Pt(1.5)):
    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, left, top, left + width, top)
    conn.line.color.rgb = color
    conn.line.width = weight
    return conn


def add_v_line(slide, left, top, height, color=LINE_RED, weight=Pt(1.5)):
    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, left, top, left, top + height)
    conn.line.color.rgb = color
    conn.line.width = weight
    return conn


def page_header(slide, page_num, total, section_title):
    add_rect(slide, Inches(0), Inches(0), Inches(0.35), H, fill_color=SEMI_DEEP)
    add_rect(slide, Inches(0.12), Inches(2.5), Inches(0.11), Inches(0.5), fill_color=BG_WHITE)

    add_text(slide, Inches(0.7), Inches(0.35), Inches(10), Inches(0.5), section_title,
             font_size=24, color=SEMI_DEEP, bold=True)
    add_h_line(slide, Inches(0.7), Inches(0.92), Inches(12), color=LINE_LIGHT, weight=Pt(1))

    add_text(slide, Inches(11.8), Inches(0.35), Inches(1.2), Inches(0.35),
             f"{page_num} / {total}", font_size=10, color=TEXT_MUTED,
             alignment=PP_ALIGN.RIGHT)

    add_h_line(slide, Inches(0.7), Inches(7.1), Inches(12), color=LINE_LIGHT, weight=Pt(0.5))
    add_text(slide, Inches(0.7), Inches(7.15), Inches(8), Inches(0.25),
             "湖南省半导体行业协会  |  2026年度",
             font_size=8, color=TEXT_MUTED)
    add_text(slide, Inches(9), Inches(7.15), Inches(3.5), Inches(0.25),
             "基于公开信息整理",
             font_size=8, color=TEXT_MUTED, alignment=PP_ALIGN.RIGHT)


TOTAL_PAGES = 7


# ========== 幻灯片 1: 封面 ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_WHITE)

add_rect(slide, Inches(0), Inches(0), Inches(4.5), H, fill_color=SEMI_DEEP)

add_rect(slide, Inches(0.6), Inches(2.8), Inches(1.2), Inches(0.03), fill_color=BG_WHITE)
add_text(slide, Inches(0.6), Inches(1.5), Inches(3.5), Inches(0.4),
         "HUNAN SEMICONDUCTOR", font_size=12, color=RGBColor(0xE8, 0xC8, 0xC0), font_name='Arial')
add_text(slide, Inches(0.6), Inches(1.85), Inches(3.5), Inches(0.4),
         "INDUSTRY ASSOCIATION", font_size=12, color=RGBColor(0xE8, 0xC8, 0xC0), font_name='Arial')

add_text(slide, Inches(0.6), Inches(5.8), Inches(3.5), Inches(0.3),
         "2026 RESEARCH REPORT", font_size=10, color=RGBColor(0xD0, 0xA0, 0xA0), font_name='Arial')

add_text(slide, Inches(5.2), Inches(2.2), Inches(7.5), Inches(0.5),
         "行业研究报告", font_size=14, color=ACCENT_GOLD, font_name='微软雅黑')

add_text(slide, Inches(5.2), Inches(2.7), Inches(7.5), Inches(1.1),
         "湖南省半导体行业协会", font_size=42, color=SEMI_DEEP, bold=True, font_name='微软雅黑')

add_h_line(slide, Inches(5.2), Inches(3.85), Inches(3.5), color=ACCENT_GOLD, weight=Pt(2))

add_multiline_text(slide, Inches(5.2), Inches(4.1), Inches(7), Inches(0.8),
                   ["芯聚潇湘 · 智链未来",
                    "凝聚发展强大合力，推进产业集群高端高质高效发展"],
                   font_size=14, color=TEXT_BODY, line_spacing=1.5)

add_text(slide, Inches(5.2), Inches(5.8), Inches(5), Inches(0.3),
         "基于公开信息整理  |  2026年7月", font_size=10, color=TEXT_MUTED)


# ========== 幻灯片 2: 协会概况 ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_WHITE)
page_header(slide, 2, TOTAL_PAGES, "协会概况")

# 关键信息卡片
card_data = [
    ("成立时间", "2022年10月10日", "长沙中电软件园"),
    ("首任会长", "丁荣军", "中国工程院院士"),
    ("发起单位", "16家", "知名半导体企业"),
    ("会员规模", "74家", "含企业及高校"),
]

for i, (label, value, sub) in enumerate(card_data):
    x = Inches(0.7) + Inches(i * 3.05)
    y = Inches(1.3)
    add_rounded_rect(slide, x, y, Inches(2.85), Inches(1.6),
                     fill_color=BG_LIGHT, line_color=LINE_LIGHT, radius=0.05)
    add_text(slide, x + Inches(0.2), y + Inches(0.2), Inches(2.45), Inches(0.3),
             label, font_size=11, color=TEXT_MUTED)
    add_text(slide, x + Inches(0.2), y + Inches(0.55), Inches(2.45), Inches(0.45),
             value, font_size=28, color=SEMI_PRIMARY, bold=True)
    add_text(slide, x + Inches(0.2), y + Inches(1.1), Inches(2.45), Inches(0.3),
             sub, font_size=10, color=TEXT_SECONDARY)

# 协会简介
add_text(slide, Inches(0.7), Inches(3.2), Inches(3), Inches(0.35),
         "▎协会简介", font_size=16, color=SEMI_DEEP, bold=True)

intro_text = (
    "湖南省半导体行业协会是由中车半导体、景嘉微电子、国科微电子、楚微半导体、三安半导体、"
    "湘能华磊光电等16家单位共同发起，于2022年10月10日在长沙中电软件园正式成立。"
    '协会以"联合会员单位聚力核心技术攻关，打通产业上下游"为宗旨，集聚行业内领先领军企业和'
    "顶尖人才，致力于成为推进湖南省半导体产业集群集聚集约、高端高质高效发展的坚实力量和重要平台。"
    "中国工程院院士丁荣军当选首任会长，罗海辉、杨国庆、王志春任副会长，王志春兼任秘书长。"
)
add_multiline_text(slide, Inches(0.7), Inches(3.65), Inches(12), Inches(1.6),
                   [intro_text[i:i+60] for i in range(0, len(intro_text), 60)],
                   font_size=11, color=TEXT_BODY, line_spacing=1.6)

# 协会地址
add_text(slide, Inches(0.7), Inches(5.6), Inches(8), Inches(0.3),
         "📍 协会地址：湖南省长沙市高新区中电软件园  |  🏛 主管单位：湖南省工业和信息化厅",
         font_size=10, color=TEXT_SECONDARY)


# ========== 幻灯片 3: 组织架构 ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_WHITE)
page_header(slide, 3, TOTAL_PAGES, "组织架构")

# 会长
add_rounded_rect(slide, Inches(5.2), Inches(1.3), Inches(3.2), Inches(1.1),
                 fill_color=SEMI_DEEP, radius=0.08)
add_text(slide, Inches(5.2), Inches(1.4), Inches(3.2), Inches(0.35),
         "会  长", font_size=12, color=ACCENT_GOLD, bold=True, alignment=PP_ALIGN.CENTER)
add_text(slide, Inches(5.2), Inches(1.75), Inches(3.2), Inches(0.45),
         "丁荣军（中国工程院院士）", font_size=14, color=TEXT_WHITE, bold=True, alignment=PP_ALIGN.CENTER)

# 竖线连接
add_v_line(slide, Inches(6.8), Inches(2.4), Inches(0.3), color=SEMI_PRIMARY, weight=Pt(1.5))

# 副会长行
for i, (name, title) in enumerate([
    ("罗海辉", "副会长"), ("杨国庆", "副会长"), ("王志春", "副会长兼秘书长")
]):
    x = Inches(2.0) + Inches(i * 3.5)
    add_rounded_rect(slide, x, Inches(2.9), Inches(3.0), Inches(0.85),
                     fill_color=BG_LIGHT, line_color=LINE_LIGHT, radius=0.05)
    add_text(slide, x, Inches(3.0), Inches(3.0), Inches(0.3),
             title, font_size=10, color=TEXT_MUTED, alignment=PP_ALIGN.CENTER)
    add_text(slide, x, Inches(3.3), Inches(3.0), Inches(0.35),
             name, font_size=16, color=SEMI_DEEP, bold=True, alignment=PP_ALIGN.CENTER)

# 三条竖线连接
add_v_line(slide, Inches(3.5), Inches(3.75), Inches(0.3), color=SEMI_PRIMARY, weight=Pt(1))
add_v_line(slide, Inches(7.0), Inches(3.75), Inches(0.3), color=SEMI_PRIMARY, weight=Pt(1))
add_v_line(slide, Inches(10.5), Inches(3.75), Inches(0.3), color=SEMI_PRIMARY, weight=Pt(1))

# 理事会 / 秘书处 / 团体标准技术委员会
org_units = [
    ("理事会", "审议人事调整、财务报告\n及新入会会员申请等事宜",
     SEMI_PRIMARY),
    ("秘书处", "负责协会日常运营管理\n协调会员服务与对外联络",
     ACCENT_TEAL),
    ("团体标准技术委员会", "制定发布团体标准\n如《SiP通用技术要求》等",
     ACCENT_GOLD),
]
for i, (name, desc, color) in enumerate(org_units):
    x = Inches(1.0) + Inches(i * 4.2)
    add_rounded_rect(slide, x, Inches(4.25), Inches(3.7), Inches(1.6),
                     fill_color=BG_LIGHT, line_color=color, line_width=Pt(1.5), radius=0.05)
    # 顶部色条
    add_rect(slide, x, Inches(4.25), Inches(3.7), Inches(0.06), fill_color=color)
    add_text(slide, x + Inches(0.2), Inches(4.45), Inches(3.3), Inches(0.35),
             name, font_size=14, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    add_multiline_text(slide, x + Inches(0.2), Inches(4.85), Inches(3.3), Inches(0.9),
                       desc.split('\n'), font_size=10, color=TEXT_BODY,
                       alignment=PP_ALIGN.CENTER, line_spacing=1.5)

# 74家会员单位
add_text(slide, Inches(0.7), Inches(6.15), Inches(12), Inches(0.3),
         "▎会员单位：74家（含企业和高校）  |  涵盖芯片设计、制造、封测、设备材料等全产业链",
         font_size=10, color=TEXT_SECONDARY, alignment=PP_ALIGN.CENTER)


# ========== 幻灯片 4: 主要发起单位 ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_WHITE)
page_header(slide, 4, TOTAL_PAGES, "主要发起单位")

founding_members = [
    ("中车半导体", "功率半导体器件龙头，IGBT/SiC器件研发制造", SEMI_PRIMARY),
    ("景嘉微电子", "国产GPU领军企业，图形处理芯片设计", ACCENT_ORANGE),
    ("国科微电子", "大规模集成电路设计，存储/视频解码芯片", ACCENT_TEAL),
    ("楚微半导体", "特色工艺集成电路制造，晶圆代工服务", ACCENT_NAVY),
    ("三安半导体", "化合物半导体龙头，碳化硅衬底/器件", ACCENT_GREEN),
    ("湘能华磊光电", "LED芯片及光电半导体器件研发制造", ACCENT_PURPLE),
]

for i, (name, desc, color) in enumerate(founding_members):
    row = i // 3
    col = i % 3
    x = Inches(0.7) + Inches(col * 4.1)
    y = Inches(1.3) + Inches(row * 2.8)

    # 卡片
    add_rounded_rect(slide, x, y, Inches(3.85), Inches(2.5),
                     fill_color=BG_WHITE, line_color=LINE_LIGHT, radius=0.05)
    # 左侧色条
    add_rect(slide, x, y, Inches(0.06), Inches(2.5), fill_color=color)
    # 序号圆圈
    add_rounded_rect(slide, x + Inches(0.25), y + Inches(0.3), Inches(0.45), Inches(0.45),
                     fill_color=color, radius=0.5)
    add_text(slide, x + Inches(0.25), y + Inches(0.33), Inches(0.45), Inches(0.4),
             str(i + 1), font_size=14, color=TEXT_WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    add_text(slide, x + Inches(0.9), y + Inches(0.3), Inches(2.7), Inches(0.35),
             name, font_size=18, color=color, bold=True)
    add_text(slide, x + Inches(0.9), y + Inches(0.75), Inches(2.7), Inches(0.15),
             "————————————————", font_size=8, color=LINE_LIGHT)
    add_multiline_text(slide, x + Inches(0.25), y + Inches(1.1), Inches(3.35), Inches(1.1),
                       [desc], font_size=11, color=TEXT_BODY, line_spacing=1.5)

# 底部说明
add_text(slide, Inches(0.7), Inches(6.5), Inches(12), Inches(0.3),
         "※ 以上为部分发起单位，共16家发起单位及74家会员单位，涵盖芯片设计、制造、封测、装备材料等全产业链环节",
         font_size=9, color=TEXT_MUTED)


# ========== 幻灯片 5: 主要职能 ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_WHITE)
page_header(slide, 5, TOTAL_PAGES, "主要职能")

functions = [
    ("🏭 产业协作与创新",
     "发挥桥梁纽带作用，组织产业链上下游企业、高校、科研院所，采取“揭榜挂帅”“赛马”等方式"
     "开展“产学研用”一体化攻关和“链式”协同攻关，推动芯片制造业高端化、智能化、绿色化发展。"),
    ("📋 团体标准制定",
     "根据《湖南省半导体行业协会团体标准管理办法》，组织制定和发布团体标准（如《系统级封装(SiP)"
     "通用技术要求》等），批准相关团体标准立项，提升行业标准化水平。"),
    ("🎓 产教融合与人才培养",
     "组织校企交流，共建实习实训基地，联合技术攻关，促进产学研用深度融合，"
     "联合培养研究生和高素质复合型人才。多所高校已入选首批高校会员单位。"),
    ("🤝 行业交流与合作",
     "组织行业年会、发布产业白皮书、与其他地区半导体行业协会及产业园区签订战略合作协议，"
     "举办主题沙龙活动，策划筹办产业峰会和创新成果展等大型品牌活动。"),
    ("📊 产业调研与政策呼吁",
     "开展企业调研、报告编写、智库建设、产业活动举办、企业对接、政策呼吁等工作，"
     "向政府反映产业发展诉求，争取政策与资源支持。"),
]

for i, (title, desc) in enumerate(functions):
    y = Inches(1.25) + Inches(i * 1.2)
    color = CAT_COLORS[i % len(CAT_COLORS)]

    # 色条标记
    add_rect(slide, Inches(0.7), y, Inches(0.05), Inches(0.95), fill_color=color)
    add_text(slide, Inches(1.0), y + Inches(0.02), Inches(11.5), Inches(0.35),
             title, font_size=15, color=color, bold=True)
    add_multiline_text(slide, Inches(1.0), y + Inches(0.38), Inches(11.5), Inches(0.55),
                       [desc], font_size=10, color=TEXT_BODY, line_spacing=1.4)


# ========== 幻灯片 6: 近期重要活动 ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_WHITE)
page_header(slide, 6, TOTAL_PAGES, "近期重要活动")

events = [
    ("第一次会员大会暨成立大会",
     "2022年10月10日",
     "在长沙中电软件园召开，选举产生第一届理事会，丁荣军院士当选会长，"
     "标志着湖南省半导体产业进入组织化协同发展新阶段。"),
    ("第一届第三次理事会",
     "2023年度",
     "审议人事调整、财务报告及新入会会员申请，黄奂果任协会秘书长，"
     "进一步完善了协会的组织架构和管理机制。"),
    ("会员大会暨工作总结会",
     "2024年度",
     '“芯聚潇湘、智链未来”主题会议，总结年度工作成果，'
     '规划下一阶段重点任务，促进会员单位交流合作。'),
    ("校企合作调研交流",
     "2024-2025年",
     "湖南师范大学等高校科创调研小组赴协会调研学习，推动产教融合，"
     "共建实习实训基地，联合培养高素质半导体人才。"),
    ("团体标准立项与发布",
     "2024-2025年",
     "发布《系统级封装(SiP)通用技术要求》等团体标准，"
     "推动湖南省半导体行业标准化、规范化发展。"),
]

for i, (title, date, desc) in enumerate(events):
    y = Inches(1.2) + Inches(i * 1.2)

    # 时间轴圆点
    add_rounded_rect(slide, Inches(1.8), y + Inches(0.15), Inches(0.18), Inches(0.18),
                     fill_color=SEMI_PRIMARY, radius=0.5)
    if i < len(events) - 1:
        add_v_line(slide, Inches(1.89), y + Inches(0.33), Inches(0.87), color=LINE_LIGHT, weight=Pt(1))

    # 日期
    add_text(slide, Inches(0.7), y, Inches(1.0), Inches(0.35),
             date, font_size=10, color=ACCENT_GOLD, bold=True)

    # 标题
    add_text(slide, Inches(2.3), y, Inches(10), Inches(0.35),
             title, font_size=15, color=SEMI_DEEP, bold=True)
    # 描述
    add_multiline_text(slide, Inches(2.3), y + Inches(0.38), Inches(10), Inches(0.6),
                       [desc], font_size=10, color=TEXT_BODY, line_spacing=1.4)


# ========== 幻灯片 7: 产业意义与展望 ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_WHITE)
page_header(slide, 7, TOTAL_PAGES, "产业意义与展望")

# 左侧：产业意义
add_text(slide, Inches(0.7), Inches(1.3), Inches(6), Inches(0.4),
         "▎协会成立的重要意义", font_size=18, color=SEMI_DEEP, bold=True)
add_h_line(slide, Inches(0.7), Inches(1.75), Inches(5.5), color=LINE_RED, weight=Pt(1.5))

significances = [
    "🔗 打通产业链上下游，形成协同创新合力",
    "🏗 集聚领军企业和顶尖人才，构建产业生态",
    "🚀 加速核心技术攻关，突破\"卡脖子\"难题",
    "📈 推动湖南半导体产业集群集聚集约发展",
    "🌐 提升湖南半导体产业在全国的知名度和影响力",
    "🎯 为政府决策提供产业智库支撑",
]

for i, txt in enumerate(significances):
    y = Inches(2.0) + Inches(i * 0.6)
    add_rounded_rect(slide, Inches(0.7), y, Inches(5.5), Inches(0.5),
                     fill_color=BG_LIGHT if i % 2 == 0 else BG_WHITE,
                     line_color=LINE_LIGHT, radius=0.03)
    add_text(slide, Inches(0.9), y + Inches(0.08), Inches(5.1), Inches(0.35),
             txt, font_size=12, color=TEXT_DARK)

# 右侧：展望
add_text(slide, Inches(7.0), Inches(1.3), Inches(6), Inches(0.4),
         "▎未来展望", font_size=18, color=SEMI_DEEP, bold=True)
add_h_line(slide, Inches(7.0), Inches(1.75), Inches(5.5), color=ACCENT_GOLD, weight=Pt(1.5))

outlook_text = (
    "随着全球半导体产业格局深刻变革，湖南作为中部地区重要的半导体产业基地，"
    "正迎来前所未有的发展机遇。湖南省半导体行业协会将继续发挥桥梁纽带作用，"
    "重点推进以下工作：\n\n"
    "• 深化“产学研用”协同创新机制\n"
    "• 完善团体标准体系建设\n"
    "• 推动第三代半导体（碳化硅/氮化镓）产业布局\n"
    "• 加强功率半导体、GPU芯片等优势领域\n"
    "• 促进长株潭半导体产业一体化发展\n"
    "• 对接长三角、珠三角产业资源\n\n"
    "在丁荣军院士的引领下，协会将凝聚74家会员单位的创新合力，"
    "将湖南打造成为国内领先、具有国际影响力的半导体产业高地。"
)
add_multiline_text(slide, Inches(7.0), Inches(2.0), Inches(5.5), Inches(4.3),
                   outlook_text.split('\n'), font_size=11, color=TEXT_BODY, line_spacing=1.5)

# 底部总结果
add_rounded_rect(slide, Inches(0.7), Inches(5.8), Inches(12.0), Inches(0.55),
                 fill_color=SEMI_PALE, line_color=SEMI_PRIMARY, line_width=Pt(1), radius=0.05)
add_text(slide, Inches(0.7), Inches(5.85), Inches(12.0), Inches(0.45),
         "芯聚潇湘，智链未来 —— 湖南省半导体行业协会将持续推进产业集群高端高质高效发展",
         font_size=13, color=SEMI_DEEP, bold=True, alignment=PP_ALIGN.CENTER)


# ========== 保存 ==========
output_path = "/Users/mrlin/Desktop/qingruan/湖南省半导体行业协会.pptx"
prs.save(output_path)
print(f"✅ PPT 已保存至: {output_path}")
print(f"   共 {TOTAL_PAGES} 页幻灯片")
