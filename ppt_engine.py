#!/usr/bin/env python3
"""PPT 生成引擎 — 共享工具函数与主题系统"""

from dataclasses import dataclass, field
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR


# === 布局常量 ===
SIDEBAR_W = Inches(0.35)
CONTENT_X = Inches(0.7)
CONTENT_W = Inches(12.0)
ACCENT_BAR_W = Inches(0.06)
ACCENT_BAR_THIN = Inches(0.04)
ACCENT_BAR_TOP = Inches(0.03)
CARD_PAD = Inches(0.15)
CARD_PAD_X = Inches(0.12)
CARD_GAP_SM = Inches(0.1)
CARD_GAP_MD = Inches(0.13)
CARD_GAP_LG = Inches(0.2)
LINE_THIN = Pt(0.5)
LINE_MED = Pt(1)
LINE_THICK = Pt(1.5)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


# === 通用文字颜色常量 ===
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
TEXT_DARK = RGBColor(0x2C, 0x2C, 0x2C)
TEXT_BODY = RGBColor(0x4A, 0x4A, 0x4A)
TEXT_SECONDARY = RGBColor(0x7A, 0x7A, 0x8A)
TEXT_MUTED = RGBColor(0xAA, 0xAA, 0xB8)


@dataclass
class Theme:
    """行业主题色彩体系"""
    name: str
    industry_label: str  # 页脚使用的产业名称，如"湖南省光电融合产业全景图"

    # 主色系
    deep: RGBColor = field(default=RGBColor(0x66, 0x00, 0x99))
    primary: RGBColor = field(default=RGBColor(0x66, 0x00, 0x99))
    light: RGBColor = field(default=RGBColor(0x8B, 0x4D, 0xB8))
    pale: RGBColor = field(default=RGBColor(0xF3, 0xEF, 0xF7))
    mist: RGBColor = field(default=RGBColor(0xE8, 0xE0, 0xF0))

    # 辅助色
    accent_gold: RGBColor = field(default=RGBColor(0xC4, 0x9A, 0x2A))
    accent_red: RGBColor = field(default=RGBColor(0x9C, 0x28, 0x28))
    accent_teal: RGBColor = field(default=RGBColor(0x00, 0x6C, 0x67))
    accent_navy: RGBColor = field(default=RGBColor(0x1A, 0x2E, 0x5C))
    accent_orange: RGBColor = field(default=RGBColor(0xD4, 0x6A, 0x1A))
    accent_green: RGBColor = field(default=RGBColor(0x2E, 0x86, 0x4B))
    accent_purple: RGBColor = field(default=RGBColor(0x6C, 0x34, 0x80))

    # 背景
    bg_white: RGBColor = field(default=RGBColor(0xFF, 0xFF, 0xFF))
    bg_light: RGBColor = field(default=RGBColor(0xF5, 0xF3, 0xF8))
    bg_sidebar: RGBColor = field(default=RGBColor(0x4B, 0x2E, 0x73))

    # 线条
    line_light: RGBColor = field(default=RGBColor(0xE0, 0xDC, 0xE8))
    line_primary: RGBColor = field(default=RGBColor(0x66, 0x00, 0x99))

    # 封面英文文字颜色（左侧色条内）
    cover_en_color: RGBColor = field(default=RGBColor(0xD0, 0xB8, 0xE8))
    cover_en_bottom_color: RGBColor = field(default=RGBColor(0xB8, 0x9A, 0xD8))
    cover_label_color: RGBColor = field(default=RGBColor(0xC4, 0x9A, 0x2A))

    # 产业链上游标签副文字颜色
    chain_sub_color: RGBColor = field(default=RGBColor(0xE0, 0xD0, 0xF0))

    # 趋势页大数字颜色
    trend_num_color: RGBColor = field(default=RGBColor(0xE8, 0xE0, 0xF0))

    # 类别色列表（用于目录/标题多样化配色）
    cat_colors: list = field(default_factory=lambda: [
        RGBColor(0x66, 0x00, 0x99),
        RGBColor(0x00, 0x6C, 0x67),
        RGBColor(0xC4, 0x9A, 0x2A),
        RGBColor(0x9C, 0x28, 0x28),
        RGBColor(0x1A, 0x2E, 0x5C),
        RGBColor(0x4B, 0x2E, 0x73),
        RGBColor(0x7A, 0x52, 0x10),
    ])


# === 预置主题 ===

THEME_OPTO = Theme(
    name="optoelectronic",
    industry_label="湖南省光电融合产业全景图",
    deep=RGBColor(0x4B, 0x2E, 0x73),
    primary=RGBColor(0x66, 0x00, 0x99),
    light=RGBColor(0x8B, 0x4D, 0xB8),
    pale=RGBColor(0xF3, 0xEF, 0xF7),
    mist=RGBColor(0xE8, 0xE0, 0xF0),
    accent_gold=RGBColor(0xC4, 0x9A, 0x2A),
    accent_red=RGBColor(0x9C, 0x28, 0x28),
    accent_teal=RGBColor(0x00, 0x6C, 0x67),
    accent_navy=RGBColor(0x1A, 0x2E, 0x5C),
    bg_light=RGBColor(0xF5, 0xF3, 0xF8),
    bg_sidebar=RGBColor(0x4B, 0x2E, 0x73),
    line_light=RGBColor(0xE0, 0xDC, 0xE8),
    line_primary=RGBColor(0x66, 0x00, 0x99),
    cover_en_color=RGBColor(0xD0, 0xB8, 0xE8),
    cover_en_bottom_color=RGBColor(0xB8, 0x9A, 0xD8),
    cover_label_color=RGBColor(0xC4, 0x9A, 0x2A),
    chain_sub_color=RGBColor(0xE0, 0xD0, 0xF0),
    trend_num_color=RGBColor(0xE8, 0xE0, 0xF0),
)

THEME_CYBER = Theme(
    name="cybersecurity",
    industry_label="湖南省网络安全产业全景图",
    deep=RGBColor(0x00, 0x3D, 0x6B),
    primary=RGBColor(0x00, 0x5B, 0x96),
    light=RGBColor(0x00, 0x89, 0xCF),
    pale=RGBColor(0xEE, 0xF4, 0xFA),
    mist=RGBColor(0xDD, 0xEB, 0xF5),
    accent_gold=RGBColor(0xB8, 0x86, 0x0A),
    accent_red=RGBColor(0xC0, 0x39, 0x2B),
    accent_teal=RGBColor(0x00, 0x7C, 0x7C),
    accent_navy=RGBColor(0x1A, 0x2E, 0x5C),
    accent_orange=RGBColor(0xD4, 0x6A, 0x1A),
    accent_green=RGBColor(0x2E, 0x86, 0x4B),
    accent_purple=RGBColor(0x6C, 0x34, 0x80),
    bg_light=RGBColor(0xF4, 0xF7, 0xFB),
    bg_sidebar=RGBColor(0x00, 0x3D, 0x6B),
    line_light=RGBColor(0xDD, 0xE2, 0xEA),
    line_primary=RGBColor(0x00, 0x5B, 0x96),
    cover_en_color=RGBColor(0xB0, 0xCF, 0xE8),
    cover_en_bottom_color=RGBColor(0x90, 0xB8, 0xD8),
    cover_label_color=RGBColor(0xD4, 0x6A, 0x1A),
    chain_sub_color=RGBColor(0xD0, 0xE0, 0xF0),
    trend_num_color=RGBColor(0xDD, 0xEB, 0xF5),
)

THEME_SEMI = Theme(
    name="semiconductor",
    industry_label="湖南省半导体产业全景图",
    deep=RGBColor(0x8B, 0x1A, 0x1A),
    primary=RGBColor(0xB5, 0x2E, 0x2E),
    light=RGBColor(0xD4, 0x4A, 0x3A),
    pale=RGBColor(0xFC, 0xF5, 0xF3),
    mist=RGBColor(0xF8, 0xEA, 0xE8),
    accent_gold=RGBColor(0xC4, 0x9A, 0x2A),
    accent_red=RGBColor(0xC0, 0x39, 0x2B),
    accent_teal=RGBColor(0x00, 0x6C, 0x67),
    accent_navy=RGBColor(0x1A, 0x2E, 0x5C),
    accent_orange=RGBColor(0xD4, 0x6A, 0x1A),
    accent_green=RGBColor(0x2E, 0x86, 0x4B),
    accent_purple=RGBColor(0x6C, 0x34, 0x80),
    bg_light=RGBColor(0xFD, 0xF7, 0xF5),
    bg_sidebar=RGBColor(0x8B, 0x1A, 0x1A),
    line_light=RGBColor(0xE8, 0xDE, 0xDA),
    line_primary=RGBColor(0xB5, 0x2E, 0x2E),
    cover_en_color=RGBColor(0xE8, 0xC8, 0xC0),
    cover_en_bottom_color=RGBColor(0xD0, 0xA0, 0xA0),
    cover_label_color=RGBColor(0xC4, 0x9A, 0x2A),
    chain_sub_color=RGBColor(0xF0, 0xE0, 0xD8),
    trend_num_color=RGBColor(0xF0, 0xE0, 0xD8),
)


# === 工具函数 ===

def create_presentation():
    """创建标准尺寸的 Presentation 对象"""
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs, prs.slide_width, prs.slide_height


def set_slide_bg(slide, color=WHITE):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, fill_color=None, line_color=None, line_width=LINE_THIN):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.shadow.inherit = False
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape


def add_rounded_rect(slide, left, top, width, height, fill_color=None, line_color=None,
                     line_width=LINE_THIN, radius=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.shadow.inherit = False
    if radius is not None:
        shape.adjustments[0] = radius
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, left, top, width, height, text, font_size=12, color=TEXT_DARK,
             bold=False, alignment=PP_ALIGN.LEFT, font_name='微软雅黑', anchor=MSO_ANCHOR.TOP,
             line_spacing=1.2):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Pt(2)
    tf.margin_right = Pt(2)
    tf.margin_top = Pt(1)
    tf.margin_bottom = Pt(1)
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    p.space_after = Pt(2)
    p.line_spacing = line_spacing
    return txBox


def add_h_line(slide, left, top, width, color, weight=LINE_THICK):
    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, left, top, left + width, top)
    conn.line.color.rgb = color
    conn.line.width = weight
    return conn


def add_v_line(slide, left, top, height, color, weight=LINE_THICK):
    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, left, top, left, top + height)
    conn.line.color.rgb = color
    conn.line.width = weight
    return conn


def page_header(slide, page_num, total, section_title, theme):
    """页眉：左侧窄边导航条 + 顶部标题区 + 底部页脚"""
    sidebar_color = theme.bg_sidebar if theme.bg_sidebar else theme.deep
    add_rect(slide, Inches(0), Inches(0), SIDEBAR_W, SLIDE_H, fill_color=sidebar_color)
    add_rect(slide, Inches(0.12), Inches(2.5), Inches(0.11), Inches(0.5), fill_color=WHITE)

    add_text(slide, CONTENT_X, Inches(0.35), Inches(10), Inches(0.5), section_title,
             font_size=24, color=theme.deep, bold=True)

    add_h_line(slide, CONTENT_X, Inches(0.92), Inches(12), color=theme.line_light, weight=LINE_MED)

    add_text(slide, Inches(11.8), Inches(0.35), Inches(1.2), Inches(0.35),
             f"{page_num} / {total}", font_size=10, color=TEXT_MUTED,
             alignment=PP_ALIGN.RIGHT)

    add_h_line(slide, CONTENT_X, Inches(7.1), Inches(12), color=theme.line_light, weight=LINE_THIN)
    add_text(slide, CONTENT_X, Inches(7.15), Inches(8), Inches(0.25),
             f"{theme.industry_label}  |  2026年度",
             font_size=8, color=TEXT_MUTED)
    add_text(slide, Inches(9), Inches(7.15), Inches(3.5), Inches(0.25),
             "基于公开信息整理",
             font_size=8, color=TEXT_MUTED, alignment=PP_ALIGN.RIGHT)
