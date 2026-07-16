#!/usr/bin/env python3
"""湖南省半导体产业全景图 PPT 生成脚本 — 红金锗色系"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# === 红金色系（半导体主题）===
SEMI_DEEP = RGBColor(0x8B, 0x1A, 0x1A)        # #8B1A1A 深红
SEMI_PRIMARY = RGBColor(0xB5, 0x2E, 0x2E)      # #B52E2E 主红
SEMI_LIGHT = RGBColor(0xD4, 0x4A, 0x3A)        # 亮红
SEMI_PALE = RGBColor(0xFC, 0xF5, 0xF3)          # 极浅红底
SEMI_MIST = RGBColor(0xF8, 0xEA, 0xE8)          # 雾红

# 辅助色
ACCENT_TEAL = RGBColor(0x00, 0x6C, 0x67)        # 深青
ACCENT_GOLD = RGBColor(0xC4, 0x9A, 0x2A)        # 金
ACCENT_ORANGE = RGBColor(0xD4, 0x6A, 0x1A)      # 橙
ACCENT_NAVY = RGBColor(0x1A, 0x2E, 0x5C)        # 藏蓝
ACCENT_GREEN = RGBColor(0x2E, 0x86, 0x4B)       # 绿
ACCENT_RED = RGBColor(0xC0, 0x39, 0x2B)          # 红
ACCENT_PURPLE = RGBColor(0x6C, 0x34, 0x80)       # 紫

BG_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BG_OFFWHITE = RGBColor(0xFA, 0xFA, 0xFC)
BG_LIGHT = RGBColor(0xFD, 0xF7, 0xF5)
BG_CARD = RGBColor(0xFF, 0xFF, 0xFF)
BG_SIDEBAR = RGBColor(0x8B, 0x1A, 0x1A)

TEXT_DARK = RGBColor(0x2C, 0x2C, 0x2C)
TEXT_BODY = RGBColor(0x4A, 0x4A, 0x4A)
TEXT_SECONDARY = RGBColor(0x7A, 0x7A, 0x8A)
TEXT_MUTED = RGBColor(0xAA, 0xAA, 0xB8)
TEXT_WHITE = RGBColor(0xFF, 0xFF, 0xFF)

LINE_LIGHT = RGBColor(0xE8, 0xDE, 0xDA)
LINE_RED = RGBColor(0xB5, 0x2E, 0x2E)

CAT_COLORS = [
    SEMI_PRIMARY, ACCENT_ORANGE, ACCENT_GOLD,
    ACCENT_TEAL, ACCENT_NAVY, ACCENT_GREEN,
    ACCENT_PURPLE,
]

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
W = prs.slide_width
H = prs.slide_height


def set_slide_bg(slide, color=BG_WHITE):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, fill_color=None, line_color=None, line_width=Pt(0.5)):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.shadow.inherit = False
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape


def add_rounded_rect(slide, left, top, width, height, fill_color=None, line_color=None, line_width=Pt(0.5), radius=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.shadow.inherit = False
    if radius is not None:
        shape.adjustments[0] = radius
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, left, top, width, height, text, font_size=12, color=TEXT_DARK,
             bold=False, alignment=PP_ALIGN.LEFT, font_name='微软雅黑', anchor=MSO_ANCHOR.TOP,
             line_spacing=1.2):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Pt(2)
    tf.margin_right = Pt(2)
    tf.margin_top = Pt(1)
    tf.margin_bottom = Pt(1)
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    p.space_after = Pt(2)
    p.line_spacing = line_spacing
    return txBox


def add_h_line(slide, left, top, width, color=LINE_RED, weight=Pt(1.5)):
    from pptx.enum.shapes import MSO_CONNECTOR
    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, left, top, left + width, top)
    conn.line.color.rgb = color
    conn.line.width = weight
    return conn


def add_v_line(slide, left, top, height, color=LINE_RED, weight=Pt(1.5)):
    from pptx.enum.shapes import MSO_CONNECTOR
    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, left, top, left, top + height)
    conn.line.color.rgb = color
    conn.line.width = weight
    return conn


def page_header(slide, page_num, total, section_title):
    add_rect(slide, Inches(0), Inches(0), Inches(0.35), H, fill_color=SEMI_DEEP)
    add_rect(slide, Inches(0.12), Inches(2.5), Inches(0.11), Inches(0.5), fill_color=BG_WHITE)

    add_text(slide, Inches(0.7), Inches(0.35), Inches(10), Inches(0.5), section_title,
             font_size=24, color=SEMI_DEEP, bold=True, font_name='微软雅黑')
    add_h_line(slide, Inches(0.7), Inches(0.92), Inches(12), color=LINE_LIGHT, weight=Pt(1))

    add_text(slide, Inches(11.8), Inches(0.35), Inches(1.2), Inches(0.35),
             f"{page_num} / {total}", font_size=10, color=TEXT_MUTED,
             alignment=PP_ALIGN.RIGHT)

    add_h_line(slide, Inches(0.7), Inches(7.1), Inches(12), color=LINE_LIGHT, weight=Pt(0.5))
    add_text(slide, Inches(0.7), Inches(7.15), Inches(8), Inches(0.25),
             "湖南省半导体产业全景图  |  2026年度",
             font_size=8, color=TEXT_MUTED)
    add_text(slide, Inches(9), Inches(7.15), Inches(3.5), Inches(0.25),
             "基于公开信息整理",
             font_size=8, color=TEXT_MUTED, alignment=PP_ALIGN.RIGHT)


TOTAL_PAGES = 9


# ========== 幻灯片 1: 封面 ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_WHITE)

add_rect(slide, Inches(0), Inches(0), Inches(4.5), H, fill_color=SEMI_DEEP)

add_rect(slide, Inches(0.6), Inches(2.8), Inches(1.2), Inches(0.03), fill_color=BG_WHITE)
add_text(slide, Inches(0.6), Inches(1.5), Inches(3.5), Inches(0.4),
         "HUNAN SEMICONDUCTOR", font_size=12, color=RGBColor(0xE8, 0xC8, 0xC0),
         font_name='Arial')
add_text(slide, Inches(0.6), Inches(1.85), Inches(3.5), Inches(0.4),
         "INDUSTRY PANORAMA", font_size=12, color=RGBColor(0xE8, 0xC8, 0xC0),
         font_name='Arial')

add_text(slide, Inches(0.6), Inches(5.8), Inches(3.5), Inches(0.3),
         "2026 INDUSTRY REPORT", font_size=10, color=RGBColor(0xD0, 0xA0, 0xA0),
         font_name='Arial')

add_text(slide, Inches(5.2), Inches(2.2), Inches(7.5), Inches(0.5),
         "产业研究报告", font_size=14, color=ACCENT_GOLD, font_name='微软雅黑')

add_text(slide, Inches(5.2), Inches(2.7), Inches(7.5), Inches(1.1),
         "湖南省半导体\n产业全景图",
         font_size=40, color=SEMI_DEEP, bold=True, font_name='微软雅黑',
         line_spacing=1.3)

add_rect(slide, Inches(5.2), Inches(4.15), Inches(0.8), Inches(0.04), fill_color=SEMI_PRIMARY)

add_text(slide, Inches(5.2), Inches(4.3), Inches(7.5), Inches(0.4),
         "功率半导体\"国家队\"  ·  SiC/GaN双轮驱动  ·  2026年度产业图谱",
         font_size=12, color=TEXT_SECONDARY, font_name='微软雅黑')

kpis = [
    ("1,571.74", "亿元", "重点项目总投资"),
    ("50", "个", "重点项目"),
    ("478.91", "亿元", "预计新增营收"),
    ("6", "个", "百亿级项目"),
]
for i, (val, unit, label) in enumerate(kpis):
    x = Inches(5.2 + i * 1.85)
    add_text(slide, x, Inches(5.2), Inches(1.7), Inches(0.5), val,
             font_size=24, color=SEMI_DEEP, bold=True, font_name='Arial')
    add_text(slide, x + Inches(1.2), Inches(5.35), Inches(0.5), Inches(0.3), unit,
             font_size=11, color=TEXT_MUTED)
    add_text(slide, x, Inches(5.65), Inches(1.7), Inches(0.3), label,
             font_size=9, color=TEXT_SECONDARY)

add_h_line(slide, Inches(5.2), Inches(6.4), Inches(7), color=LINE_LIGHT, weight=Pt(0.5))
add_text(slide, Inches(5.2), Inches(6.5), Inches(7), Inches(0.3),
         "数据来源：湖南省工业和信息化厅 · 红网 · 中国证券报 · 时代电气/三安光电年报",
         font_size=8, color=TEXT_MUTED)


# ========== 幻灯片 2: 目录 ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_WHITE)
page_header(slide, 2, TOTAL_PAGES, "目  录")

contents = [
    ("01", "产业链全景", "半导体上中下游全链路解析", SEMI_DEEP),
    ("02", "功率半导体", "IGBT与SiC/GaN双赛道深度", ACCENT_ORANGE),
    ("03", "区域布局", "六大核心产业集聚区", ACCENT_TEAL),
    ("04", "重点企业", "16家龙头与骨干企业", ACCENT_GOLD),
    ("05", "创新平台", "科研平台与高校力量", ACCENT_NAVY),
    ("06", "政策支撑", "8项核心政策体系", ACCENT_RED),
]

for i, (num, title, desc, color) in enumerate(contents):
    col = i % 2
    row = i // 2
    x = Inches(0.9 + col * 6.0)
    y = Inches(1.4 + row * 1.7)

    add_text(slide, x, y, Inches(1.2), Inches(0.8), num,
             font_size=48, color=color, bold=True, font_name='Arial',
             anchor=MSO_ANCHOR.MIDDLE)
    add_v_line(slide, x + Inches(1.3), y + Inches(0.15), Inches(0.7), color=LINE_LIGHT, weight=Pt(1))
    add_text(slide, x + Inches(1.5), y + Inches(0.1), Inches(4), Inches(0.4), title,
             font_size=18, color=TEXT_DARK, bold=True, font_name='微软雅黑')
    add_text(slide, x + Inches(1.5), y + Inches(0.55), Inches(4), Inches(0.3), desc,
             font_size=11, color=TEXT_SECONDARY, font_name='微软雅黑')


# ========== 幻灯片 3: 产业链全景 ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_WHITE)
page_header(slide, 3, TOTAL_PAGES, "半导体产业链全景")

chain_data = [
    ("上游", "材料 · 设备", SEMI_DEEP, [
        ("衬底/外延", "SiC衬底、GaN外延片\n硅基外延、蓝宝石衬底", "三安半导体 · 赛德雷特"),
        ("芯片设计", "CPU/GPU/SSD控制器/DSP\nFPGA、模拟芯片、射频芯片", "景嘉微 · 国科微 · 飞腾 · 毂梁微"),
        ("半导体材料", "基板玻璃、电子特气、光刻胶\n靶材、CMP抛光液", "邵虹 · 江丰电子 · 金博股份"),
        ("制造设备", "晶体生长炉、刻蚀设备\n光刻机配件、测试设备", "阿秒光学 · 弘宇精密"),
        ("封装材料", "封装基板、引线框架\n陶瓷基板、键合丝", "信维电科 · 明正宏"),
    ]),
    ("中游", "制造 · 封测", ACCENT_ORANGE, [
        ("晶圆制造", "8英寸SiC晶圆产线（国内首批）\n6英寸SiC月产能16000片", "湖南三安 · 中车时代"),
        ("功率半导体", "IGBT/IGCT/MOSFET IDM模式\n中低压功率器件全产业链", "中车时代半导体"),
        ("先进封装", "SiP/3D封装、晶圆级封装\n功率器件封装", "国创越摩 · 湘潭智造基地"),
        ("第三代半导体", "SiC MOSFET/GaN HEMT\n8英寸SiC产线规模化通线", "三安半导体 · 棣山科技"),
        ("存储器/逻辑", "长鑫存储生态协同\n国产存储产业链配套", "岳阳紫光 · 国科微"),
    ]),
    ("下游", "应用 · 终端", ACCENT_GOLD, [
        ("新能源汽车", "车规级SiC主驱芯片\n理想/小鹏/蔚来主驱供应", "三安半导体 · 中车时代"),
        ("轨道交通", "高铁/地铁IGBT模块\n大功率晶闸管全球第一", "中车时代半导体"),
        ("智能电网", "特高压功率器件\n柔性直流输电换流阀", "中车时代 · 特变电工"),
        ("AI服务器", "SiC电源芯片、GPU算力\n数据中心供电系统", "三安 · 景嘉微 · 图灵智算"),
        ("消费电子", "手机快充GaN器件\nLED/Mini-LED直显", "三安 · 惠科 · 蓝思科技"),
    ]),
]

y_start = Inches(1.15)
row_h = Inches(1.88)
for row_idx, (label, sub, color, items) in enumerate(chain_data):
    y = y_start + row_idx * row_h
    add_rect(slide, Inches(0.7), y, Inches(1.3), Inches(1.65), fill_color=color)
    add_text(slide, Inches(0.7), y + Inches(0.3), Inches(1.3), Inches(0.5), label,
             font_size=20, color=TEXT_WHITE, bold=True, alignment=PP_ALIGN.CENTER,
             font_name='微软雅黑')
    add_text(slide, Inches(0.7), y + Inches(0.85), Inches(1.3), Inches(0.3), sub,
             font_size=9, color=RGBColor(0xF0, 0xE0, 0xD8), alignment=PP_ALIGN.CENTER,
             font_name='微软雅黑')

    card_w = Inches(2.22)
    card_gap = Inches(0.08)
    for i, (title, desc, ent) in enumerate(items):
        cx = Inches(2.15) + i * (card_w + card_gap)
        add_rect(slide, cx, y, card_w, Inches(1.65), fill_color=BG_LIGHT,
                 line_color=LINE_LIGHT, line_width=Pt(0.5))
        add_rect(slide, cx, y, card_w, Inches(0.03), fill_color=color)
        add_text(slide, cx + Inches(0.1), y + Inches(0.15), card_w - Inches(0.2), Inches(0.3),
                 title, font_size=12, color=color, bold=True, font_name='微软雅黑')
        add_text(slide, cx + Inches(0.1), y + Inches(0.5), card_w - Inches(0.2), Inches(0.6),
                 desc, font_size=9, color=TEXT_BODY, font_name='微软雅黑', line_spacing=1.3)
        add_h_line(slide, cx + Inches(0.1), y + Inches(1.15), card_w - Inches(0.2),
                   color=LINE_LIGHT, weight=Pt(0.5))
        add_text(slide, cx + Inches(0.1), y + Inches(1.22), card_w - Inches(0.2), Inches(0.35),
                 ent, font_size=8, color=color, font_name='微软雅黑')


# ========== 幻灯片 4: 功率半导体深度 ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_WHITE)
page_header(slide, 4, TOTAL_PAGES, "功率半导体：IGBT与SiC/GaN双赛道")

# 左：IGBT硅基赛道
add_rect(slide, Inches(0.7), Inches(1.2), Inches(5.8), Inches(5.5),
         fill_color=BG_LIGHT, line_color=LINE_LIGHT, line_width=Pt(0.5))
add_rect(slide, Inches(0.7), Inches(1.2), Inches(5.8), Inches(0.06), fill_color=SEMI_PRIMARY)

add_text(slide, Inches(1.0), Inches(1.4), Inches(5.3), Inches(0.4),
         "硅基IGBT — 中车时代半导体", font_size=20, color=SEMI_DEEP, bold=True, font_name='微软雅黑')

igbt_items = [
    ("产业规模", "株洲功率半导体集群规模595亿元，220+企业，11个国家级创新平台"),
    ("技术地位", "全球少数同时掌握IGBT/IGCT/SiC的IDM企业，大功率晶闸管全球市占率第一"),
    ("2025营收", "半导体板块53.6亿元（+30.43%），器件销量近700万只"),
    ("产能建设", "中低压功率器件产业化（株洲）三期，宜兴三期同步推进"),
    ("下游应用", "高铁/地铁牵引、新能源汽车电驱、特高压电网、风电变流器"),
    ("技术路线", "第四代沟槽栅IGBT已突破，精细平面栅批量交付"),
]
for i, (k, v) in enumerate(igbt_items):
    y = Inches(2.0 + i * 0.72)
    add_text(slide, Inches(1.0), y, Inches(1.2), Inches(0.3), k,
             font_size=10, color=SEMI_PRIMARY, bold=True, font_name='微软雅黑')
    add_text(slide, Inches(2.3), y, Inches(4.0), Inches(0.6), v,
             font_size=10, color=TEXT_BODY, font_name='微软雅黑', line_spacing=1.3)

# 右：SiC/GaN赛道
add_rect(slide, Inches(6.7), Inches(1.2), Inches(5.8), Inches(5.5),
         fill_color=BG_LIGHT, line_color=LINE_LIGHT, line_width=Pt(0.5))
add_rect(slide, Inches(6.7), Inches(1.2), Inches(5.8), Inches(0.06), fill_color=ACCENT_ORANGE)

add_text(slide, Inches(7.0), Inches(1.4), Inches(5.3), Inches(0.4),
         "SiC/GaN第三代 — 三安半导体", font_size=20, color=ACCENT_ORANGE, bold=True, font_name='微软雅黑')

sic_items = [
    ("投资规模", "长沙产业园总投资160亿元，1000亩，国内首条8英寸SiC全产业链贯通平台"),
    ("产能现状", "6英寸SiC月产能16000片，8英寸衬底1000片/月、外延2000片/月；GaN 2000片/月"),
    ("2025营收", "湖南三安9.1亿元，2026Q1出货量同比+123%"),
    ("车规突破", "SiC MOSFET上车理想汽车高压平台，车规级主驱芯片规模化装车"),
    ("客户阵容", "理想、台达、维谛、光宝、长城、伟创力等全球头部客户批量供货"),
    ("远期规划", "达产年产36万片6英寸+48万片8英寸SiC晶圆，累计出货超3亿颗"),
]
for i, (k, v) in enumerate(sic_items):
    y = Inches(2.0 + i * 0.72)
    add_text(slide, Inches(7.0), y, Inches(1.2), Inches(0.3), k,
             font_size=10, color=ACCENT_ORANGE, bold=True, font_name='微软雅黑')
    add_text(slide, Inches(8.3), y, Inches(4.0), Inches(0.6), v,
             font_size=10, color=TEXT_BODY, font_name='微软雅黑', line_spacing=1.3)


# ========== 幻灯片 5: 区域布局 ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_WHITE)
page_header(slide, 5, TOTAL_PAGES, "六大核心产业集聚区")

regions = [
    ("长沙", "第三代半导体+芯片设计", "SiC/GaN晶圆制造 · GPU/SSD/CPU设计", "16个项目·118.78亿", SEMI_DEEP,
     "三安半导体 · 景嘉微 · 国科微 · 飞腾 · 湘江鲲鹏 · 图灵智算 · 中国长城"),
    ("株洲", "功率半导体\"国家队\"", "IGBT全产业链 · 国家级产业集群", "集群规模595亿", ACCENT_ORANGE,
     "中车时代半导体 · 时代电气 · 时代新材 · 三一硅能 · 赛德雷特"),
    ("娄底", "半导体材料新势力", "电子陶瓷 · 显示材料 · 芯片半导体", "21.16亿·+27.1%", ACCENT_GREEN,
     "美程陶瓷 · 安地亚斯 · 华菱安赛乐米塔尔 · 20家规上企业"),
    ("邵阳", "基板玻璃基地", "显示基板玻璃 · 上下游配套集聚", "4条热端产线", ACCENT_NAVY,
     "邵虹基板玻璃 · 致成科技"),
    ("湘潭", "先进封装+制造基地", "半导体制造 · 封装测试 · 智能终端", "23.8亿投资", ACCENT_TEAL,
     "半导体湘潭智造基地 · 国创越摩先进封装 · 蓝思智能终端"),
    ("岳阳", "存储器+中部基地", "新紫光集团中部基地 · 存储产业配套", "百亿级项目", ACCENT_PURPLE,
     "紫光集团 · 国科微存储 · 岳阳临港半导体产业园"),
]

card_w = Inches(4.0)
card_h = Inches(2.65)
gap_x = Inches(0.13)
gap_y = Inches(0.18)
start_x = Inches(0.7)
start_y = Inches(1.15)

for i, (city, badge, role, data_str, color, enterprises) in enumerate(regions):
    col = i % 3
    row = i // 3
    x = start_x + col * (card_w + gap_x)
    y = start_y + row * (card_h + gap_y)

    add_rect(slide, x, y, card_w, card_h, fill_color=BG_WHITE,
             line_color=LINE_LIGHT, line_width=Pt(0.5))
    add_rect(slide, x, y, Inches(0.06), card_h, fill_color=color)
    add_text(slide, x + Inches(0.2), y + Inches(0.12), Inches(2.5), Inches(0.25), badge,
             font_size=9, color=color, bold=True, font_name='微软雅黑')
    add_text(slide, x + Inches(0.2), y + Inches(0.38), Inches(2), Inches(0.45), city,
             font_size=26, color=TEXT_DARK, bold=True, font_name='微软雅黑')
    add_text(slide, x + Inches(2.5), y + Inches(0.45), card_w - Inches(2.6), Inches(0.35), data_str,
             font_size=13, color=color, bold=True, alignment=PP_ALIGN.RIGHT, font_name='微软雅黑')
    add_text(slide, x + Inches(0.2), y + Inches(0.95), card_w - Inches(0.3), Inches(0.3), role,
             font_size=10, color=TEXT_SECONDARY, font_name='微软雅黑')
    add_h_line(slide, x + Inches(0.2), y + Inches(1.3), card_w - Inches(0.4),
               color=LINE_LIGHT, weight=Pt(0.5))
    add_text(slide, x + Inches(0.2), y + Inches(1.4), card_w - Inches(0.3), Inches(1.1), enterprises,
             font_size=9, color=TEXT_BODY, font_name='微软雅黑', line_spacing=1.5)


# ========== 幻灯片 6: 重点企业 ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_WHITE)
page_header(slide, 6, TOTAL_PAGES, "龙头与重点企业")

enterprises = [
    ("中车时代半导体", "株洲", "功率半导体IDM龙头，2025半导体收入53.6亿(+30%)，全球少数IGBT+SiC全技术企业", "国家队", ACCENT_RED),
    ("三安半导体", "长沙·高新区", "国内SiC产业先行者，160亿产业园，8英寸SiC全产业链贯通，上车理想汽车", "SiC龙头", ACCENT_RED),
    ("景嘉微", "长沙", "国产GPU龙头，军民融合，\"两芯一生态\"GPU芯片唯一供应商", "芯片设计", ACCENT_TEAL),
    ("国科微", "长沙", "国产SSD控制器芯片龙头，\"七大类芯片\"之一，国家级专精特新小巨人", "芯片设计", ACCENT_TEAL),
    ("飞腾信息", "长沙(研发中心)", "国产CPU龙头之一，飞腾CPU为\"两芯一生态\"核心，信创CPU主力", "\"两芯\"核心", ACCENT_ORANGE),
    ("邵虹基板玻璃", "邵阳", "显示基板玻璃国产替代，3条热端+2条冷端产线，第4条热端在设计阶段", "十大产业项目", ACCENT_RED),
    ("金博股份", "益阳", "碳基材料龙头，半导体热场材料，创新联合体，碳化硅涂层技术领先", "创新联合体", ACCENT_GOLD),
    ("湘江鲲鹏", "长沙", "华为鲲鹏生态核心伙伴(拓维信息控股90%)，国产服务器/PC整机", "鲲鹏生态", ACCENT_GREEN),
    ("国创越摩", "湘潭", "先进封装项目(二期)，SiP/3D封装技术，填补湖南封测产业链关键环节", "先进封装", ACCENT_NAVY),
    ("图灵智算", "长沙", "宽谱域光电+量子计算，AI算力与半导体协同", "量子+AI", ACCENT_PURPLE),
    ("江丰电子", "益阳", "全球靶材龙头，益阳基地全球最大外埠，半导体溅射靶材年产值超7.5亿", "靶材龙头", ACCENT_TEAL),
    ("蓝思科技", "长沙·浏阳", "3D玻璃+智能装备，年产1万台自动化设备+50万台机器人，半导体应用材料", "十大产业项目", ACCENT_RED),
    ("阿秒光学", "郴州", "激光蚀刻设备，半导体制造精密加工，良品率98.6%，上市后备企业", "精密设备", ACCENT_GOLD),
    ("信维电科", "益阳", "MLCC多层陶瓷电容器，5G/半导体配套高端电子元器件", "电子元器件", ACCENT_GREEN),
    ("美程陶瓷", "娄底", "新能源汽车电子陶瓷，半导体封装用陶瓷基板，娄底半导体材料龙头", "陶瓷基板", ACCENT_ORANGE),
    ("赛德雷特", "株洲", "半导体衬底材料，SiC/GaN外延配套，株洲功率半导体集群关键企业", "衬底材料", ACCENT_NAVY),
]

card_w = Inches(3.05)
card_h = Inches(1.35)
gap_x = Inches(0.1)
gap_y = Inches(0.1)
start_x = Inches(0.7)
start_y = Inches(1.1)

for i, (name, loc, desc, badge, badge_color) in enumerate(enterprises):
    col = i % 4
    row = i // 4
    x = start_x + col * (card_w + gap_x)
    y = start_y + row * (card_h + gap_y)

    add_rect(slide, x, y, card_w, card_h, fill_color=BG_WHITE,
             line_color=LINE_LIGHT, line_width=Pt(0.5))
    add_rect(slide, x, y, Inches(0.04), card_h, fill_color=badge_color if badge else LINE_LIGHT)
    add_text(slide, x + Inches(0.12), y + Inches(0.06), card_w - Inches(0.2), Inches(0.3),
             name, font_size=13, color=TEXT_DARK, bold=True, font_name='微软雅黑')
    add_text(slide, x + Inches(0.12), y + Inches(0.36), card_w - Inches(0.2), Inches(0.2),
             loc, font_size=8, color=TEXT_MUTED, font_name='微软雅黑')
    add_text(slide, x + Inches(0.12), y + Inches(0.58), card_w - Inches(0.2), Inches(0.55),
             desc, font_size=8, color=TEXT_BODY, font_name='微软雅黑', line_spacing=1.3)
    if badge:
        add_text(slide, x + Inches(0.12), y + Inches(1.1), Inches(1.2), Inches(0.2), badge,
                 font_size=8, color=badge_color, bold=True, font_name='微软雅黑')


# ========== 幻灯片 7: 创新平台 ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_WHITE)
page_header(slide, 7, TOTAL_PAGES, "创新平台与科研支撑")

platforms = [
    ("功率半导体行业联盟", "株洲·国家级行业组织\n涵盖材料-芯片-模块-应用全产业链协同创新", SEMI_DEEP),
    ("湖南省集成电路产业联盟", "长沙经开区2025年成立\n整合国科微、国防科大等资源，配套50亿元专项基金", ACCENT_ORANGE),
    ("国家功率半导体创新中心", "株洲·国家级创新平台\nIGBT/SiC技术研发+中试+检测认证", ACCENT_GOLD),
    ("湖南三安半导体研究院", "长沙高新区\nSiC/GaN材料与器件研发，8英寸产线技术攻关", ACCENT_TEAL),
    ("高校科研力量", "国防科大 · 湖南大学 · 中南大学\n在功率器件、封装技术、EDA工具等领域深度布局", ACCENT_NAVY),
    ("第三代半导体产业研究院", "长沙·产学研协同\n聚焦SiC/GaN材料、器件、模组产业化技术", ACCENT_GREEN),
    ("湖南省半导体行业协会", "全省行业组织\n政策对接、产业研究、企业服务、招商引资", ACCENT_PURPLE),
    ("功率半导体可靠性实验室", "株洲中车+高校联合\n功率器件可靠性测试、寿命评估、失效分析", ACCENT_RED),
    ("集成电路产业园（长沙）", "长沙经开区200亩\n配套50亿元专项基金，吸引第三代半导体企业入驻", ACCENT_ORANGE),
]

card_w = Inches(4.0)
card_h = Inches(1.68)
gap_x = Inches(0.13)
gap_y = Inches(0.15)
start_x = Inches(0.7)
start_y = Inches(1.15)

for i, (title, desc, color) in enumerate(platforms):
    col = i % 3
    row = i // 3
    x = start_x + col * (card_w + gap_x)
    y = start_y + row * (card_h + gap_y)

    add_rect(slide, x, y, card_w, card_h, fill_color=BG_LIGHT,
             line_color=LINE_LIGHT, line_width=Pt(0.5))
    add_rounded_rect(slide, x + Inches(0.15), y + Inches(0.15), Inches(0.4), Inches(0.4),
                    fill_color=color, radius=0.5)
    add_text(slide, x + Inches(0.15), y + Inches(0.15), Inches(0.4), Inches(0.4), str(i+1),
             font_size=14, color=TEXT_WHITE, bold=True, alignment=PP_ALIGN.CENTER,
             font_name='Arial', anchor=MSO_ANCHOR.MIDDLE)
    add_text(slide, x + Inches(0.65), y + Inches(0.15), card_w - Inches(0.8), Inches(0.4), title,
             font_size=12, color=TEXT_DARK, bold=True, font_name='微软雅黑',
             anchor=MSO_ANCHOR.MIDDLE)
    add_text(slide, x + Inches(0.15), y + Inches(0.7), card_w - Inches(0.3), Inches(0.9), desc,
             font_size=9, color=TEXT_BODY, font_name='微软雅黑', line_spacing=1.4)


# ========== 幻灯片 8: 政策支撑 ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_WHITE)
page_header(slide, 8, TOTAL_PAGES, "政策支撑体系")

policies = [
    ("顶层规划", "湖南省\"十五五\"电子信息产业规划", "谋划半导体产业布局，功率半导体+第三代半导体列入重点方向", SEMI_DEEP),
    ("重点项目", "2026电子信息制造业重点项目", "总投资1571.74亿元，50个项目，6个百亿级，聚焦功率半导体全产业链", ACCENT_ORANGE),
    ("产业集群", "国家级功率半导体产业集群", "株洲功率半导体集群入选国家队，集聚220+企业、11个国家级创新平台", ACCENT_GOLD),
    ("专项基金", "集成电路50亿专项基金", "长沙经开区配套50亿元专项基金，重点吸引第三代半导体、先进封装", ACCENT_TEAL),
    ("园区载体", "长沙集成电路产业园200亩", "整合国科微+国防科大资源，2025年成立集成电路产业联盟", ACCENT_GREEN),
    ("人才政策", "半导体人才产教融合", "国防科大/湖大/中南大学半导体相关学科，领军团队最高1亿项目支持", ACCENT_NAVY),
    ("税收优惠", "集成电路企业所得税减免", "\"十免\"\"五免五减半\"等税收优惠，研发费用加计扣除", ACCENT_PURPLE),
    ("成果转化", "\"先用后付\"科技成果转化", "首批537项科技成果推广，支持半导体共性技术加速产业化", ACCENT_RED),
]

card_w = Inches(6.1)
card_h = Inches(1.32)
gap_x = Inches(0.2)
gap_y = Inches(0.1)
start_x = Inches(0.7)
start_y = Inches(1.15)

for i, (tag, title, desc, color) in enumerate(policies):
    col = i % 2
    row = i // 2
    x = start_x + col * (card_w + gap_x)
    y = start_y + row * (card_h + gap_y)

    add_rect(slide, x, y, card_w, card_h, fill_color=BG_WHITE,
             line_color=LINE_LIGHT, line_width=Pt(0.5))
    add_rect(slide, x, y, Inches(0.06), card_h, fill_color=color)
    add_text(slide, x + Inches(0.2), y + Inches(0.1), Inches(0.8), Inches(0.25), tag,
             font_size=9, color=color, bold=True, font_name='微软雅黑')
    add_text(slide, x + Inches(1.1), y + Inches(0.08), card_w - Inches(1.2), Inches(0.35), title,
             font_size=13, color=TEXT_DARK, bold=True, font_name='微软雅黑',
             anchor=MSO_ANCHOR.MIDDLE)
    add_text(slide, x + Inches(0.2), y + Inches(0.55), card_w - Inches(0.3), Inches(0.65), desc,
             font_size=10, color=TEXT_BODY, font_name='微软雅黑', line_spacing=1.3)


# ========== 幻灯片 9: 发展趋势 ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_WHITE)

add_rect(slide, Inches(0), Inches(0), Inches(0.35), H, fill_color=SEMI_DEEP)

add_text(slide, Inches(0.8), Inches(0.8), Inches(10), Inches(0.5),
         "总结与展望", font_size=28, color=SEMI_DEEP, bold=True, font_name='微软雅黑')
add_h_line(slide, Inches(0.8), Inches(1.4), Inches(11.8), color=LINE_LIGHT, weight=Pt(1))

add_text(slide, Inches(0.8), Inches(1.7), Inches(11.5), Inches(0.8),
         "湖南半导体产业 — 从\"功率半导体\"到\"第三代半导体\"双轮驱动",
         font_size=24, color=TEXT_DARK, bold=True, font_name='微软雅黑',
         alignment=PP_ALIGN.CENTER)

# 四大趋势
trends = [
    ("01", "SiC/GaN进入放量期", SEMI_DEEP,
     "湖南三安8英寸SiC全产业链贯通、中车时代8英寸SiC产线通线，2026年均进入产能爬坡期。车规级SiC主驱芯片规模化装车，AI服务器电源SiC需求爆发"),
    ("02", "功率半导体全链自主", ACCENT_ORANGE,
     "中车时代IGBT/SiC双技术路线并进，第四代沟槽栅突破。株洲功率半导体集群规模595亿，从\"卡脖子\"到全球市占率第一，正朝\"世界前三强\"目标迈进"),
    ("03", "先进封装补链强链", ACCENT_TEAL,
     "国创越摩先进封装二期、湘潭智造基地推进，SiP/3D封装填补产业链关键环节。SiC模块封装、银烧结等新技术加速导入"),
    ("04", "产业生态加速成型", ACCENT_NAVY,
     "2026年50个重点项目、总投资超1500亿。长沙设计+株洲制造+湘潭封测+娄底材料+邵阳基板+岳阳存储，六大集聚区差异化协同"),
]

card_w = Inches(2.95)
card_h = Inches(3.8)
gap = Inches(0.15)
start_x = Inches(0.7)
start_y = Inches(2.8)

for i, (num, title, color, text) in enumerate(trends):
    x = start_x + i * (card_w + gap)
    add_rect(slide, x, start_y, card_w, card_h, fill_color=BG_LIGHT,
             line_color=LINE_LIGHT, line_width=Pt(0.5))
    add_rect(slide, x, start_y, card_w, Inches(0.06), fill_color=color)
    add_text(slide, x + Inches(0.2), start_y + Inches(0.25), Inches(2), Inches(1.0), num,
             font_size=48, color=RGBColor(0xF0, 0xE0, 0xD8), bold=True, font_name='Arial')
    add_rect(slide, x + Inches(0.2), start_y + Inches(1.3), Inches(0.6), Inches(0.03),
             fill_color=color)
    add_text(slide, x + Inches(0.2), start_y + Inches(1.5), card_w - Inches(0.4), Inches(0.45), title,
             font_size=16, color=color, bold=True, font_name='微软雅黑')
    add_text(slide, x + Inches(0.2), start_y + Inches(2.1), card_w - Inches(0.4), Inches(1.5), text,
             font_size=10, color=TEXT_BODY, font_name='微软雅黑', line_spacing=1.6)

# 底部
add_h_line(slide, Inches(0.8), Inches(6.7), Inches(11.5), color=LINE_LIGHT, weight=Pt(0.5))
add_text(slide, Inches(0.8), Inches(6.75), Inches(11.5), Inches(0.3),
         "数据来源：湖南省工业和信息化厅 · 红网 · 中国证券报 · 时代电气年报 · 三安光电年报  |  2026年度",
         font_size=9, color=TEXT_MUTED, alignment=PP_ALIGN.CENTER, font_name='微软雅黑')


# === 保存 ===
output_path = "/Users/mrlin/Desktop/qingruan/湖南省半导体产业全景图.pptx"
prs.save(output_path)
print(f"PPT 已保存至: {output_path}")
print(f"共 {len(prs.slides)} 页幻灯片")
