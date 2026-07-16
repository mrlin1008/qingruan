#!/usr/bin/env python3
"""湖南省光电融合产业全景图 PPT 生成脚本 — 清华风格"""

from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

from ppt_engine import (
    THEME_OPTO, SLIDE_H, WHITE, TEXT_DARK, TEXT_BODY, TEXT_SECONDARY, TEXT_MUTED,
    CONTENT_X, LINE_THIN, LINE_MED, LINE_THICK,
    create_presentation, set_slide_bg, add_rect, add_rounded_rect, add_text,
    add_h_line, add_v_line, page_header,
)

# === 清华色彩别名 ===
THU_PURPLE = THEME_OPTO.primary
THU_PURPLE_DEEP = THEME_OPTO.deep
THU_PURPLE_MIST = THEME_OPTO.mist
ACCENT_GOLD = THEME_OPTO.accent_gold
ACCENT_RED = THEME_OPTO.accent_red
ACCENT_TEAL = THEME_OPTO.accent_teal
ACCENT_NAVY = THEME_OPTO.accent_navy
BG_WHITE = WHITE
BG_LIGHT = THEME_OPTO.bg_light
TEXT_WHITE = WHITE
LINE_LIGHT = THEME_OPTO.line_light
LINE_PURPLE = THEME_OPTO.line_primary

prs, W, H = create_presentation()


TOTAL_PAGES = 9


# ========== 幻灯片 1: 封面 ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

# 左侧紫色装饰区（留边风格）
add_rect(slide, Inches(0), Inches(0), Inches(4.5), H, fill_color=THU_PURPLE)

# 紫色区内的装饰线
add_rect(slide, Inches(0.6), Inches(2.8), Inches(1.2), Inches(0.03), fill_color=BG_WHITE)
# 紫色区内的英文文字
add_text(slide, Inches(0.6), Inches(1.5), Inches(3.5), Inches(0.4),
         "HUNAN OPTOELECTRONIC", font_size=12, color=RGBColor(0xD0, 0xB8, 0xE8),
         font_name='Arial')
add_text(slide, Inches(0.6), Inches(1.85), Inches(3.5), Inches(0.4),
         "FUSION INDUSTRY PANORAMA", font_size=12, color=RGBColor(0xD0, 0xB8, 0xE8),
         font_name='Arial')

# 紫色区底部
add_text(slide, Inches(0.6), Inches(5.8), Inches(3.5), Inches(0.3),
         "2026 INDUSTRY REPORT", font_size=10, color=RGBColor(0xB8, 0x9A, 0xD8),
         font_name='Arial')

# 右侧白色区 — 主标题
add_text(slide, Inches(5.2), Inches(2.2), Inches(7.5), Inches(0.5),
         "产业研究报告", font_size=14, color=ACCENT_GOLD, font_name='微软雅黑')

add_text(slide, Inches(5.2), Inches(2.7), Inches(7.5), Inches(1.1),
         "湖南省光电融合\n产业全景图",
         font_size=40, color=THU_PURPLE, bold=True, font_name='微软雅黑',
         line_spacing=1.3)

# 紫色短线装饰
add_rect(slide, Inches(5.2), Inches(4.15), Inches(0.8), Inches(0.04), fill_color=THU_PURPLE)

# 副标题
add_text(slide, Inches(5.2), Inches(4.3), Inches(7.5), Inches(0.4),
         "4×4 现代化产业体系  ·  十五五规划关键之年  ·  2026年度产业图谱",
         font_size=12, color=TEXT_SECONDARY, font_name='微软雅黑')

# KPI 简要数据
kpis = [
    ("1,571.74", "亿元", "重点项目总投资"),
    ("50", "个", "重点项目"),
    ("478.91", "亿元", "预计新增营收"),
    ("7", "大", "核心集聚城市"),
]
for i, (val, unit, label) in enumerate(kpis):
    x = Inches(5.2 + i * 1.85)
    add_text(slide, x, Inches(5.2), Inches(1.7), Inches(0.5), val,
             font_size=24, color=THU_PURPLE, bold=True, font_name='Arial')
    add_text(slide, x + Inches(1.2), Inches(5.35), Inches(0.5), Inches(0.3), unit,
             font_size=11, color=TEXT_MUTED)
    add_text(slide, x, Inches(5.65), Inches(1.7), Inches(0.3), label,
             font_size=9, color=TEXT_SECONDARY)

# 底部信息
add_h_line(slide, Inches(5.2), Inches(6.4), Inches(7), color=LINE_LIGHT, weight=Pt(0.5))
add_text(slide, Inches(5.2), Inches(6.5), Inches(7), Inches(0.3),
         "数据来源：湖南省工业和信息化厅 · 湖南省人民政府门户网站 · 湖南日报",
         font_size=8, color=TEXT_MUTED)


# ========== 幻灯片 2: 目录 ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
page_header(slide, 2, TOTAL_PAGES, "目  录", THEME_OPTO)

contents = [
    ("01", "产业链全景", "光电融合上中下游全链路解析", THU_PURPLE),
    ("02", "区域布局", "七大核心产业集聚区", ACCENT_TEAL),
    ("03", "重点企业", "16家龙头与重点企业", ACCENT_GOLD),
    ("04", "创新平台", "科研平台与高校力量", THU_PURPLE_DEEP),
    ("05", "政策支撑", "8项核心政策体系", ACCENT_RED),
    ("06", "发展趋势", "四大未来方向", ACCENT_NAVY),
]

for i, (num, title, desc, color) in enumerate(contents):
    col = i % 2
    row = i // 2
    x = Inches(0.9 + col * 6.0)
    y = Inches(1.4 + row * 1.7)

    # 编号
    add_text(slide, x, y, Inches(1.2), Inches(0.8), num,
             font_size=48, color=color, bold=True, font_name='Arial',
             anchor=MSO_ANCHOR.MIDDLE)
    # 竖线
    add_v_line(slide, x + Inches(1.3), y + Inches(0.15), Inches(0.7), color=LINE_LIGHT, weight=Pt(1))
    # 标题
    add_text(slide, x + Inches(1.5), y + Inches(0.1), Inches(4), Inches(0.4), title,
             font_size=18, color=TEXT_DARK, bold=True, font_name='微软雅黑')
    # 描述
    add_text(slide, x + Inches(1.5), y + Inches(0.55), Inches(4), Inches(0.3), desc,
             font_size=11, color=TEXT_SECONDARY, font_name='微软雅黑')


# ========== 幻灯片 3: 产业链全景 ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
page_header(slide, 3, TOTAL_PAGES, "光电融合产业链全景", THEME_OPTO)

chain_data = [
    ("上游", "材料 · 芯片", THU_PURPLE, [
        ("光电材料", "基板玻璃、光电特种气体、靶材、光纤预制棒", "邵虹 · 中化蓝天 · 江丰电子"),
        ("光芯片", "DFB/EML激光芯片、探测器、硅光芯片", "硅基Micro-LED · 高速EML"),
        ("功率半导体", "IGBT、MOSFET、第三代半导体", "中车时代半导体 · 三一硅能"),
        ("光纤光缆", "单模/多模光纤、特种光纤、MPO连接器", "信维电子 · 艾迪奥"),
        ("光学元器件", "光学镜头、偏光片、光学透镜、滤光片", "山嘉光电 · 谱特光电"),
    ]),
    ("中游", "器件 · 模组", ACCENT_TEAL, [
        ("新型显示面板", "LCD/OLED面板、Mini-LED背光/直显模组", "惠科 · 中沛光电 · 蓝思科技"),
        ("光通信器件", "光模块、光收发组件、激光器、光放大器", "光智通信 · 图灵智算"),
        ("光电传感与精密", "激光蚀刻设备、光电传感器、COG绑定", "阿秒光学 · 弘宇精密"),
        ("半导体封测", "功率器件封装、集成电路制造、智能装备", "湘潭智造基地 · 明正宏"),
        ("光电功能材料", "碳基材料、光电薄膜、3D玻璃、柔性材料", "金博股份 · 麓邦光电"),
    ]),
    ("下游", "应用 · 终端", ACCENT_GOLD, [
        ("智能终端", "智能手机、智能穿戴、具身智能机器人", "蓝思机器人 · 中沛手机"),
        ("数据中心与算力", "AI智算中心、超算中心、高速光互联", "图灵智算"),
        ("5G通信与网络", "5G基站光模块、光传输网络、运营商集采", "光智通信 · 信维电科"),
        ("安防与车载光电", "安防监控光电系统、车载激光雷达", "英飞拓"),
        ("新型显示终端", "大尺寸显示整机、LED直显、Mini-LED背光", "惠科直显 · 明和数艺"),
    ]),
]

y_start = Inches(1.15)
row_h = Inches(1.88)
for row_idx, (label, sub, color, items) in enumerate(chain_data):
    y = y_start + row_idx * row_h
    # 左侧标签
    add_rect(slide, Inches(0.7), y, Inches(1.3), Inches(1.65), fill_color=color)
    add_text(slide, Inches(0.7), y + Inches(0.3), Inches(1.3), Inches(0.5), label,
             font_size=20, color=TEXT_WHITE, bold=True, alignment=PP_ALIGN.CENTER,
             font_name='微软雅黑')
    add_text(slide, Inches(0.7), y + Inches(0.85), Inches(1.3), Inches(0.3), sub,
             font_size=9, color=RGBColor(0xE0, 0xD0, 0xF0), alignment=PP_ALIGN.CENTER,
             font_name='微软雅黑')

    # 5个卡片
    card_w = Inches(2.22)
    card_gap = Inches(0.08)
    for i, (title, desc, ent) in enumerate(items):
        cx = Inches(2.15) + i * (card_w + card_gap)
        add_rect(slide, cx, y, card_w, Inches(1.65), fill_color=BG_LIGHT,
                 line_color=LINE_LIGHT, line_width=Pt(0.5))
        # 顶部色条
        add_rect(slide, cx, y, card_w, Inches(0.03), fill_color=color)
        # 标题
        add_text(slide, cx + Inches(0.1), y + Inches(0.15), card_w - Inches(0.2), Inches(0.3),
                 title, font_size=12, color=color, bold=True, font_name='微软雅黑')
        # 描述
        add_text(slide, cx + Inches(0.1), y + Inches(0.5), card_w - Inches(0.2), Inches(0.6),
                 desc, font_size=9, color=TEXT_BODY, font_name='微软雅黑', line_spacing=1.3)
        # 底部企业
        add_h_line(slide, cx + Inches(0.1), y + Inches(1.15), card_w - Inches(0.2),
                   color=LINE_LIGHT, weight=Pt(0.5))
        add_text(slide, cx + Inches(0.1), y + Inches(1.22), card_w - Inches(0.2), Inches(0.35),
                 ent, font_size=8, color=color, font_name='微软雅黑')


# ========== 幻灯片 4: 七大核心集聚区 ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
page_header(slide, 4, TOTAL_PAGES, "七大核心产业集聚区", THEME_OPTO)

regions = [
    ("长沙", "省会 · 核心引擎", "新型显示龙头 · 智能终端 · 先进计算", "百亿+项目6个", THU_PURPLE,
     "惠科 · 蓝思科技 · 图灵智算 · 麓邦光电 · 明和数艺 · 韶光芯材"),
    ("株洲", "功率半导体国家队", "功率半导体全产业链 · 国家级产业集群", "220+企业", ACCENT_TEAL,
     "中车时代半导体 · 三一硅能 · 赛德雷特"),
    ("湘潭", "半导体智造基地", "半导体制造 · 智能终端设备", "23.8亿投资", ACCENT_GOLD,
     "半导体湘潭智造基地 · 蓝思智能终端 · 锦智光电 · 金杯电工"),
    ("益阳", "湖南光电谷", "光通信器件 · 光学镜头 · 长益科创走廊节点", "86项专利 · 12项转化", THU_PURPLE_DEEP,
     "未来光电技术研究院 · 金博股份 · 信维电科 · 光智通信 · 江丰电子 · 明正宏"),
    ("郴州", "湘南光电谷", "光电显示完整产业链 · 湾区产业转移承接地", "80亿产值 · 59家企业", ACCENT_RED,
     "宜章经开区 · 北湖湘南光电产业园 · 山嘉光电 · 谱特光电 · 中沛光电 · 英飞拓 · 阿秒光学"),
    ("邵阳", "基板玻璃基地", "显示基板玻璃 · 上下游配套集聚", "4条热端生产线", ACCENT_NAVY,
     "邵虹基板玻璃 · 致成科技"),
]

card_w = Inches(4.0)
card_h = Inches(2.65)
gap_x = Inches(0.13)
gap_y = Inches(0.18)
start_x = Inches(0.7)
start_y = Inches(1.15)

for i, (city, badge, role, data_str, color, enterprises) in enumerate(regions):
    col = i % 3
    row = i // 3
    x = start_x + col * (card_w + gap_x)
    y = start_y + row * (card_h + gap_y)

    add_rect(slide, x, y, card_w, card_h, fill_color=BG_WHITE,
             line_color=LINE_LIGHT, line_width=Pt(0.5))
    # 左侧色条
    add_rect(slide, x, y, Inches(0.06), card_h, fill_color=color)
    # badge
    add_text(slide, x + Inches(0.2), y + Inches(0.12), Inches(2), Inches(0.25), badge,
             font_size=9, color=color, bold=True, font_name='微软雅黑')
    # 城市名
    add_text(slide, x + Inches(0.2), y + Inches(0.38), Inches(2), Inches(0.45), city,
             font_size=26, color=TEXT_DARK, bold=True, font_name='微软雅黑')
    # 数据
    add_text(slide, x + Inches(2.5), y + Inches(0.45), card_w - Inches(2.6), Inches(0.35), data_str,
             font_size=13, color=color, bold=True, alignment=PP_ALIGN.RIGHT, font_name='微软雅黑')
    # 角色
    add_text(slide, x + Inches(0.2), y + Inches(0.95), card_w - Inches(0.3), Inches(0.3), role,
             font_size=10, color=TEXT_SECONDARY, font_name='微软雅黑')
    # 分隔线
    add_h_line(slide, x + Inches(0.2), y + Inches(1.3), card_w - Inches(0.4),
               color=LINE_LIGHT, weight=Pt(0.5))
    # 企业列表
    add_text(slide, x + Inches(0.2), y + Inches(1.4), card_w - Inches(0.3), Inches(1.1), enterprises,
             font_size=9, color=TEXT_BODY, font_name='微软雅黑', line_spacing=1.5)


# ========== 幻灯片 5: 重点企业 ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
page_header(slide, 5, TOTAL_PAGES, "龙头与重点企业", THEME_OPTO)

enterprises = [
    ("蓝思科技", "长沙·浏阳", "智能装备基地，年产1万台自动化设备+50万台机器人，3D玻璃", "十大产业项目", ACCENT_RED),
    ("惠科", "长沙·浏阳", "Mini-LED背光/直显模组及整机，主厂房封顶，百亿级项目", "十大产业项目", ACCENT_RED),
    ("中车时代半导体", "株洲", "中低压功率器件通线，国家级功率半导体集群链主", "十大产业项目", ACCENT_RED),
    ("邵虹", "邵阳", "基板玻璃，3条热端+2条冷端生产线，第4条设计阶段", "十大产业项目", ACCENT_RED),
    ("金博股份", "益阳", "光电材料龙头，创新联合体，碳基材料技术领先", "创新联合体", ACCENT_GOLD),
    ("信维电科", "益阳", "电子元器件、5G配套，高端电子元器件核心企业", "", ACCENT_TEAL),
    ("光智通信", "益阳", "光通信器件研发制造，长益常科创走廊关键企业", "", ACCENT_TEAL),
    ("江丰电子", "益阳", "全球靶材龙头，益阳基地全球最大外埠，年产值超7.5亿", "", ACCENT_TEAL),
    ("山嘉光电", "郴州·宜章", "大尺寸偏光片供应商，服务惠科/华星/京东方", "", THU_PURPLE),
    ("谱特光电", "郴州·宜章", "打破偏光片日韩垄断，超薄圆偏光片入选省首套件", "首套件", ACCENT_GOLD),
    ("中沛光电", "郴州·北湖", "链主型企业，5秒/台手机整机，COG技术填补空白", "", THU_PURPLE),
    ("英飞拓", "郴州·北湖", "全国安防5强，年产值预计突破3亿元", "", THU_PURPLE),
    ("图灵智算", "长沙", "宽谱域高端光电产品，量子科技+光电融合", "", ACCENT_NAVY),
    ("麓邦光电", "长沙", "显示屏项目，光电功能材料与新型显示", "", ACCENT_NAVY),
    ("明正宏电子", "益阳", "双层/多层线路板扩建，年产能增至150万平方米", "", ACCENT_TEAL),
    ("阿秒光学", "郴州·北湖", "激光蚀刻设备，与中沛形成技术闭环，良品率98.6%", "上市后备", ACCENT_GOLD),
]

card_w = Inches(3.05)
card_h = Inches(1.35)
gap_x = Inches(0.1)
gap_y = Inches(0.1)
start_x = Inches(0.7)
start_y = Inches(1.1)

for i, (name, loc, desc, badge, badge_color) in enumerate(enterprises):
    col = i % 4
    row = i // 4
    x = start_x + col * (card_w + gap_x)
    y = start_y + row * (card_h + gap_y)

    add_rect(slide, x, y, card_w, card_h, fill_color=BG_WHITE,
             line_color=LINE_LIGHT, line_width=Pt(0.5))
    # 左色条
    add_rect(slide, x, y, Inches(0.04), card_h, fill_color=badge_color if badge else LINE_LIGHT)
    # 名称
    add_text(slide, x + Inches(0.12), y + Inches(0.06), card_w - Inches(0.2), Inches(0.3),
             name, font_size=13, color=TEXT_DARK, bold=True, font_name='微软雅黑')
    # 地点
    add_text(slide, x + Inches(0.12), y + Inches(0.36), card_w - Inches(0.2), Inches(0.2),
             loc, font_size=8, color=TEXT_MUTED, font_name='微软雅黑')
    # 描述
    add_text(slide, x + Inches(0.12), y + Inches(0.58), card_w - Inches(0.2), Inches(0.55),
             desc, font_size=8, color=TEXT_BODY, font_name='微软雅黑', line_spacing=1.3)
    # badge
    if badge:
        add_text(slide, x + Inches(0.12), y + Inches(1.1), Inches(1.2), Inches(0.2), badge,
                 font_size=8, color=badge_color, bold=True, font_name='微软雅黑')


# ========== 幻灯片 6: 创新平台 ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
page_header(slide, 6, TOTAL_PAGES, "创新平台与科研支撑", THEME_OPTO)

platforms = [
    ("湖南未来光电技术研究院", "益阳高新区×湖南师范大学共建\n突破硅基Micro-LED，申报专利86项，成果转化12项", THU_PURPLE),
    ("5G+电容器科技孵化器", "益阳高新区\n科技型企业孵化器，聚焦光电产业创新孵化", ACCENT_TEAL),
    ("新型电子元器件中试基地", "益阳高新区\n科技成果转化中试基地，支撑光电元器件产业化", ACCENT_GOLD),
    ("高校科研力量", "国防科大 · 湖南大学 · 中南大学\n湖南师大 · 长沙理工 · 湘南学院", THU_PURPLE_DEEP),
    ("金博股份创新联合体", "光电材料领域\n产学研协同攻关创新联合体", ACCENT_RED),
    ("国家新型显示产业联盟", "宜章经开区光电产业协会加入\n拓宽产业发展空间", ACCENT_NAVY),
    ("新型显示研发中心", "宜章经开区投资390万元筹建\n产学研深度合作，4名国家级专家", ACCENT_TEAL),
    ("郴江实验室", "郴州地区科研平台\n帮助企业解决技术痛点、研发难点", THU_PURPLE),
    ("飞地园区 · 科创飞地", "益阳高新区-湘江新区飞地园区\n益阳（长沙）光电技术创新中心", THU_PURPLE_DEEP),
]

card_w = Inches(4.0)
card_h = Inches(1.68)
gap_x = Inches(0.13)
gap_y = Inches(0.15)
start_x = Inches(0.7)
start_y = Inches(1.15)

for i, (title, desc, color) in enumerate(platforms):
    col = i % 3
    row = i // 3
    x = start_x + col * (card_w + gap_x)
    y = start_y + row * (card_h + gap_y)

    add_rect(slide, x, y, card_w, card_h, fill_color=BG_LIGHT,
             line_color=LINE_LIGHT, line_width=Pt(0.5))
    # 编号圆
    add_rounded_rect(slide, x + Inches(0.15), y + Inches(0.15), Inches(0.4), Inches(0.4),
                    fill_color=color, radius=0.5)
    add_text(slide, x + Inches(0.15), y + Inches(0.15), Inches(0.4), Inches(0.4), str(i+1),
             font_size=14, color=TEXT_WHITE, bold=True, alignment=PP_ALIGN.CENTER,
             font_name='Arial', anchor=MSO_ANCHOR.MIDDLE)
    # 标题
    add_text(slide, x + Inches(0.65), y + Inches(0.15), card_w - Inches(0.8), Inches(0.4), title,
             font_size=12, color=TEXT_DARK, bold=True, font_name='微软雅黑',
             anchor=MSO_ANCHOR.MIDDLE)
    # 描述
    add_text(slide, x + Inches(0.15), y + Inches(0.7), card_w - Inches(0.3), Inches(0.9), desc,
             font_size=9, color=TEXT_BODY, font_name='微软雅黑', line_spacing=1.4)


# ========== 幻灯片 7: 政策支撑 ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
page_header(slide, 7, TOTAL_PAGES, "政策支撑体系", THEME_OPTO)

policies = [
    ("顶层规划", "湖南省\"十五五\"未来产业发展规划", "谋划光电产业发展布局，光电信息列入重点方向", THU_PURPLE),
    ("产业体系", "\"4×4\"现代化产业体系", "功率半导体纳入重点培育，聚焦新型显示、人工智能", ACCENT_TEAL),
    ("区域协同", "\"长益\"未来光电技术科创与产业走廊", "长沙-益阳科创走廊，规划衔接、政策协同、产业联动", ACCENT_GOLD),
    ("专项支持", "支持益阳高端电子元器件产业", "重点发展激光器、光传感器、硅光芯片，打造湖南光电谷", THU_PURPLE_DEEP),
    ("成果转化", "\"先用后付\"推广实施方案", "首批推广537项科技成果，支持光电共性技术加速转化", ACCENT_RED),
    ("飞地机制", "飞地园区发展若干措施", "支持GDP核算、税收分成、能耗指标分配", ACCENT_NAVY),
    ("知识产权", "知识产权服务集聚试点", "益阳高新区纳入试点，一站式综合服务", THU_PURPLE),
    ("区域规划", "益阳\"十五五\"新型工业化规划", "光电产业为战略性主导产业，明确湖南光电谷定位", THU_PURPLE_DEEP),
]

card_w = Inches(6.1)
card_h = Inches(1.32)
gap_x = Inches(0.2)
gap_y = Inches(0.1)
start_x = Inches(0.7)
start_y = Inches(1.15)

for i, (tag, title, desc, color) in enumerate(policies):
    col = i % 2
    row = i // 2
    x = start_x + col * (card_w + gap_x)
    y = start_y + row * (card_h + gap_y)

    add_rect(slide, x, y, card_w, card_h, fill_color=BG_WHITE,
             line_color=LINE_LIGHT, line_width=Pt(0.5))
    # 左色条
    add_rect(slide, x, y, Inches(0.06), card_h, fill_color=color)
    # tag
    add_text(slide, x + Inches(0.2), y + Inches(0.1), Inches(0.8), Inches(0.25), tag,
             font_size=9, color=color, bold=True, font_name='微软雅黑')
    # title
    add_text(slide, x + Inches(1.1), y + Inches(0.08), card_w - Inches(1.2), Inches(0.35), title,
             font_size=13, color=TEXT_DARK, bold=True, font_name='微软雅黑',
             anchor=MSO_ANCHOR.MIDDLE)
    # desc
    add_text(slide, x + Inches(0.2), y + Inches(0.55), card_w - Inches(0.3), Inches(0.65), desc,
             font_size=10, color=TEXT_BODY, font_name='微软雅黑', line_spacing=1.3)


# ========== 幻灯片 8: 发展趋势 ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
page_header(slide, 8, TOTAL_PAGES, "发展趋势与未来方向", THEME_OPTO)

trends = [
    ("01", "AI + 光电融合", THU_PURPLE,
     "AI算力大规模运算需求驱动光电技术战略地位提升，光互联替代电互联成为算力网络核心支撑，硅光CPO技术加速落地"),
    ("02", "国产替代加速", ACCENT_TEAL,
     "高端光芯片、EML、高速DSP国产化率不足10%，湖南在功率半导体、靶材、偏光片等领域实现突破，\"从0到1\"持续涌现"),
    ("03", "长益科创走廊", ACCENT_GOLD,
     "长沙研发+益阳转化模式深化，飞地园区机制打通区域壁垒，长益常科创走廊光电产业带加速成型"),
    ("04", "湾区产业承接", ACCENT_RED,
     "湖南\"一带一部\"区位优势凸显，50个重点项目中24%承接长三角、粤港澳大湾区产业转移，郴州成为核心承接地"),
]

card_w = Inches(2.95)
card_h = Inches(4.8)
gap = Inches(0.15)
start_x = Inches(0.7)
start_y = Inches(1.3)

for i, (num, title, color, text) in enumerate(trends):
    x = start_x + i * (card_w + gap)
    add_rect(slide, x, start_y, card_w, card_h, fill_color=BG_LIGHT,
             line_color=LINE_LIGHT, line_width=Pt(0.5))
    # 顶部色块
    add_rect(slide, x, start_y, card_w, Inches(0.06), fill_color=color)
    # 大数字
    add_text(slide, x + Inches(0.2), start_y + Inches(0.25), Inches(2), Inches(1.2), num,
             font_size=56, color=THU_PURPLE_MIST, bold=True, font_name='Arial')
    # 横线
    add_rect(slide, x + Inches(0.2), start_y + Inches(1.55), Inches(0.6), Inches(0.03),
             fill_color=color)
    # 标题
    add_text(slide, x + Inches(0.2), start_y + Inches(1.75), card_w - Inches(0.4), Inches(0.5), title,
             font_size=18, color=color, bold=True, font_name='微软雅黑')
    # 描述
    add_text(slide, x + Inches(0.2), start_y + Inches(2.5), card_w - Inches(0.4), Inches(2), text,
             font_size=10, color=TEXT_BODY, font_name='微软雅黑', line_spacing=1.6)


# ========== 幻灯片 9: 总结 ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

# 左侧紫色装饰
add_rect(slide, Inches(0), Inches(0), Inches(0.35), H, fill_color=THU_PURPLE)

# 标题
add_text(slide, Inches(0.8), Inches(0.8), Inches(10), Inches(0.5),
         "总结与展望", font_size=28, color=THU_PURPLE, bold=True, font_name='微软雅黑')
add_h_line(slide, Inches(0.8), Inches(1.4), Inches(11.8), color=LINE_LIGHT, weight=Pt(1))

# 主标题
add_text(slide, Inches(0.8), Inches(1.7), Inches(11.5), Inches(0.8),
         "湖南光电融合产业 — 迈向新质生产力高地",
         font_size=24, color=TEXT_DARK, bold=True, font_name='微软雅黑',
         alignment=PP_ALIGN.CENTER)

summary_items = [
    ("产业规模", "1,571.74亿元重点项目总投资，50个重点项目，6个百亿级项目", THU_PURPLE),
    ("区域格局", "长沙-株洲-湘潭-益阳-郴州-邵阳-衡阳，七大集聚区协同发展", ACCENT_TEAL),
    ("创新突破", "硅基Micro-LED、功率半导体、偏光片国产替代，\"从0到1\"持续涌现", ACCENT_GOLD),
    ("政策护航", "\"十五五\"规划引领、长益科创走廊、飞地园区、先用后付成果转化", ACCENT_RED),
    ("未来方向", "AI+光电融合、国产替代加速、湾区产业承接、长益常科创走廊成型", ACCENT_NAVY),
]

for i, (title, desc, color) in enumerate(summary_items):
    y = Inches(2.9 + i * 0.72)
    add_rect(slide, Inches(1.0), y, Inches(11), Inches(0.6), fill_color=BG_LIGHT,
             line_color=LINE_LIGHT, line_width=Pt(0.5))
    add_rect(slide, Inches(1.0), y, Inches(0.06), Inches(0.6), fill_color=color)
    add_text(slide, Inches(1.2), y, Inches(2), Inches(0.6), title,
             font_size=14, color=color, bold=True, font_name='微软雅黑',
             anchor=MSO_ANCHOR.MIDDLE)
    add_text(slide, Inches(3.3), y, Inches(8.5), Inches(0.6), desc,
             font_size=11, color=TEXT_BODY, font_name='微软雅黑',
             anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.4)

# 底部
add_h_line(slide, Inches(0.8), Inches(6.6), Inches(11.5), color=LINE_LIGHT, weight=Pt(0.5))
add_text(slide, Inches(0.8), Inches(6.7), Inches(11.5), Inches(0.3),
         "数据来源：湖南省工业和信息化厅 · 湖南省人民政府门户网站 · 湖南日报  |  2026年度",
         font_size=9, color=TEXT_MUTED, alignment=PP_ALIGN.CENTER, font_name='微软雅黑')


# === 保存 ===
output_path = "/Users/mrlin/Desktop/qingruan/湖南省光电融合产业全景图.pptx"
prs.save(output_path)
print(f"PPT 已保存至: {output_path}")
print(f"共 {len(prs.slides)} 页幻灯片")
